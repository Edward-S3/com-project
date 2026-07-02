"""各種ファイル形式からのテキスト抽出"""
from __future__ import annotations

import io
from pathlib import Path

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".md", ".docx", ".pptx", ".xlsx"}


class DocumentParseError(Exception):
    """ファイル解析に失敗した場合"""


def extract_text(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise DocumentParseError(f"未対応の形式です: {ext}")

    handlers = {
        ".txt": _extract_txt,
        ".md": _extract_txt,
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
    }
    try:
        text = handlers[ext](data)
    except DocumentParseError:
        raise
    except Exception as exc:
        raise DocumentParseError(f"{filename} の読み込みに失敗しました: {exc}") from exc

    cleaned = (text or "").strip()
    if not cleaned:
        raise DocumentParseError(f"{filename} からテキストを抽出できませんでした。")
    return cleaned


def extract_multiple(files: list[tuple[str, bytes]]) -> str:
    parts: list[str] = []
    for name, data in files:
        body = extract_text(name, data)
        parts.append(f"### ファイル: {name}\n{body}")
    return "\n\n".join(parts)


def _extract_txt(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp932", "shift_jis", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    raise DocumentParseError("テキストの文字コードを判定できませんでした。")


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages: list[str] = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(pages)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"[スライド {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in wb.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(",".join(cells))
            if rows:
                parts.append(f"[シート: {sheet.title}]\n" + "\n".join(rows))
    finally:
        wb.close()
    return "\n\n".join(parts)
