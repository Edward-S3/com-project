"""V5: JSON → pptx（GridCell / CHAPTER_TITLE / Unicodeアイコン描画）"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.config import FONT_BODY, FONT_TITLE, SLIDE_MARGIN_IN
from src.schemas import GridCell, PresentationData, SlideContent, SlideTheme, VisualElement
from src.utils.font_sizer import calculate_font_size

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(SLIDE_MARGIN_IN)
CONTENT_W = SLIDE_WIDTH - MARGIN * 2
GRID_CELL_PAD_IN = 0.22
GRID_LIGHT_BG = "#EBF0FF"
GRID_SEP_LIGHT = "#C8D8F8"
EMOJI_ICON_PT = 32
STEP_FLOW_EMOJIS = ("💡", "🔍", "⚙️", "📈", "✅")
ALERT_CHAPTER_KEYWORDS = ("警告", "リスク", "セキュリティ", "注意", "危険", "コンプライアンス", "脅威")


def _normalize_hex(hex_color: str, *, fallback: str = "FFFFFF") -> str:
    value = (hex_color or f"#{fallback}").strip().lstrip("#")
    if len(value) == 3:
        value = "".join(c * 2 for c in value)
    if len(value) != 6:
        value = fallback
    return value.upper()


def hex_to_rgb(hex_color: str) -> RGBColor:
    value = _normalize_hex(hex_color)
    r, g, b = (int(value[i : i + 2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


def _relative_luminance(hex_color: str) -> float:
    r, g, b = (int(_normalize_hex(hex_color)[i : i + 2], 16) for i in (0, 2, 4))

    def _linear(channel: int) -> float:
        c = channel / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * _linear(r) + 0.7152 * _linear(g) + 0.0722 * _linear(b)


def _contrast_ratio(hex_a: str, hex_b: str) -> float:
    l1 = _relative_luminance(hex_a)
    l2 = _relative_luminance(hex_b)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _pick_readable_text(bg_hex: str, *candidates: str) -> str:
    options = [c for c in candidates if c]
    if not options:
        return "#000000"
    return max(options, key=lambda color: _contrast_ratio(bg_hex, color))


def _blend_hex(c1: str, c2: str, ratio: float) -> str:
    a, b = hex_to_rgb(c1), hex_to_rgb(c2)
    r = int(a[0] + (b[0] - a[0]) * ratio)
    g = int(a[1] + (b[1] - a[1]) * ratio)
    bl = int(a[2] + (b[2] - a[2]) * ratio)
    return f"#{r:02X}{g:02X}{bl:02X}"


def _emu_inches(length) -> float:
    return float(length) / 914400 if isinstance(length, (int, float)) else length.inches


def _cell_bg(theme: SlideTheme, variant: str) -> str:
    mapping = {
        "dominant": theme.dominant_color,
        "support": theme.support_color,
        "accent": theme.accent_color,
        "light": GRID_LIGHT_BG,
    }
    return mapping.get(variant, theme.support_color)


def _grid_cell_colors(bg: str, theme: SlideTheme) -> tuple[str, str]:
    """GRID セルの前景色と区切り線色を返す (fg, separator)"""
    bg_norm = _normalize_hex(bg)
    dom_norm = _normalize_hex(theme.dominant_color)
    is_dark = bg_norm == dom_norm or _relative_luminance(bg) < 0.45
    fg = "#FFFFFF" if is_dark else theme.dominant_color
    sep_col = "#FFFFFF" if is_dark else GRID_SEP_LIGHT
    return fg, sep_col


def draw_grid_cell(
    slide,
    cell: GridCell,
    left,
    top,
    width,
    height,
    theme: SlideTheme,
) -> None:
    """GRID レイアウトの1セル: 背景 + header + 区切り線 + body を個別描画"""
    bg = _cell_bg(theme, cell.color_variant)
    fg, sep_col = _grid_cell_colors(bg, theme)

    l_in = _emu_inches(left)
    t_in = _emu_inches(top)
    w_in = _emu_inches(width)
    h_in = _emu_inches(height)
    pad = GRID_CELL_PAD_IN
    inner_w = w_in - pad * 2

    # 背景
    bg_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(l_in),
        Inches(t_in),
        Inches(w_in),
        Inches(h_in),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = hex_to_rgb(bg)
    bg_shape.line.fill.background()

    # 見出し (header)
    header_h_in = h_in * 0.36
    tb_h = slide.shapes.add_textbox(
        Inches(l_in + pad),
        Inches(t_in + pad),
        Inches(inner_w),
        Inches(header_h_in),
    )
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    tf_h.vertical_anchor = MSO_ANCHOR.TOP
    tf_h.margin_left = tf_h.margin_right = tf_h.margin_top = tf_h.margin_bottom = Pt(0)
    ph = tf_h.paragraphs[0]
    ph.text = cell.header
    ph.alignment = PP_ALIGN.LEFT
    ph.font.name = FONT_TITLE
    ph.font.bold = True
    ph.font.size = Pt(15)
    ph.font.color.rgb = hex_to_rgb(fg)

    # 区切り線
    sep_top_in = t_in + h_in * 0.38
    sep = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(l_in + pad),
        Inches(sep_top_in - 0.04),
        Inches(inner_w),
        Inches(0.015),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = hex_to_rgb(sep_col)
    sep.line.fill.background()

    # 本文 (body)
    body_top_in = sep_top_in + 0.02
    body_h_in = h_in - h_in * 0.40 - 0.08
    tb_b = slide.shapes.add_textbox(
        Inches(l_in + pad),
        Inches(body_top_in),
        Inches(inner_w),
        Inches(body_h_in),
    )
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    tf_b.vertical_anchor = MSO_ANCHOR.TOP
    tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Pt(0)
    body_lines = cell.body[:3]
    row_h_in = body_h_in / max(len(body_lines), 1)
    for i, line in enumerate(body_lines):
        p = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
        p.text = f"・{line}"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        pt = calculate_font_size(line, inner_w, row_h_in, 11)
        p.font.name = FONT_BODY
        p.font.size = Pt(pt)
        p.font.color.rgb = hex_to_rgb(fg)


def build_pptx(presentation: PresentationData) -> bytes:
    prs = Presentation()
    prs.slide_width = int(SLIDE_WIDTH)
    prs.slide_height = int(SLIDE_HEIGHT)
    theme = presentation.theme
    total = len(presentation.slides)

    for idx, slide_data in enumerate(presentation.slides):
        is_first = idx == 0
        is_last = idx == total - 1
        layout = slide_data.layout_type

        if layout == "TITLE":
            _render_title(prs, slide_data, theme, dark=is_first or (is_last and theme.is_dark_mode_cover))
        elif layout == "CHAPTER_TITLE":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            draw_chapter_title_slide(slide, slide_data, theme)
        elif layout == "GRID_2X2":
            _render_grid(prs, slide_data, theme, cols=2)
        elif layout == "GRID_1X3":
            _render_grid(prs, slide_data, theme, cols=3)
        elif layout == "FULL_STEP_FLOW":
            _render_full_step_flow(prs, slide_data, theme)
        elif layout == "MATURITY_BAR_FULL":
            _render_maturity_full(prs, slide_data, theme)
        else:
            _render_two_column(prs, slide_data, theme)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _is_alert_chapter(content: SlideContent) -> bool:
    text = f"{content.title} {content.key_message} {content.subtitle or ''}"
    return any(keyword in text for keyword in ALERT_CHAPTER_KEYWORDS)


def draw_chapter_title_slide(slide, content: SlideContent, theme: SlideTheme) -> None:
    """章扉スライド: 中央配置。警告章は alert_color、通常は dominant_color"""
    bg = theme.alert_color if _is_alert_chapter(content) else theme.dominant_color
    text_color = "#FFFFFF"
    _set_bg(slide, bg)
    _add_text(
        slide, MARGIN, Inches(2.2), CONTENT_W, Inches(1.4),
        content.title, 36, True, text_color, PP_ALIGN.CENTER,
    )
    sub = content.subtitle or content.key_message
    if sub:
        _add_text(
            slide, MARGIN, Inches(3.8), CONTENT_W, Inches(1.0),
            sub, 20, False, text_color, PP_ALIGN.CENTER,
        )


def _add_emoji_icon(
    slide,
    left,
    top,
    width,
    height,
    emoji: str,
    *,
    color_hex: str,
    size_pt: int = EMOJI_ICON_PT,
) -> None:
    """Unicode 絵文字を大きなフォントで描画"""
    icon = (emoji or "•").strip()
    if not icon:
        icon = "•"
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = icon[:4]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(size_pt)
    p.font.color.rgb = hex_to_rgb(color_hex)


def _extract_step_emoji(item: str, index: int) -> tuple[str, str]:
    """STEP_FLOW 項目から先頭絵文字とラベル本文を分離"""
    text = (item or "").strip()
    if text and ord(text[0]) > 0x2600:
        parts = text.split(" ", 1)
        if len(parts) == 1:
            return parts[0], ""
        return parts[0], parts[1]
    return STEP_FLOW_EMOJIS[index % len(STEP_FLOW_EMOJIS)], text


def _set_bg(slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def _add_rect(slide, left, top, width, height, fill_hex: str, *, rounded: bool = True) -> None:
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shape.line.fill.background()


def _render_header(slide, slide_data: SlideContent, theme: SlideTheme) -> None:
    _set_bg(slide, "#FFFFFF")
    _add_text(slide, MARGIN, MARGIN, CONTENT_W, Inches(0.55), slide_data.title, 22, True, theme.dominant_color)
    if slide_data.key_message:
        km_color = _pick_readable_text("#FFFFFF", theme.accent_color, theme.dominant_color)
        _add_text(slide, MARGIN, Inches(1.05), CONTENT_W, Inches(0.35), slide_data.key_message, 14, False, km_color)


def _add_text(slide, left, top, width, height, text, base_pt, bold, color_hex, align=PP_ALIGN.LEFT) -> None:
    pt = calculate_font_size(text, _emu_inches(width), _emu_inches(height), base_pt)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    font = p.font
    font.name = FONT_TITLE if bold else FONT_BODY
    font.size = Pt(pt)
    font.bold = bold
    font.color.rgb = hex_to_rgb(color_hex)


def _render_title(prs, slide_data: SlideContent, theme: SlideTheme, *, dark: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = theme.dominant_color if dark else "#FFFFFF"
    text_color = "#FFFFFF" if dark else theme.dominant_color
    _set_bg(slide, bg)
    _add_text(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(1.2), slide_data.title, 40, True, text_color)
    sub = slide_data.subtitle or slide_data.key_message or (slide_data.bullet_points[0] if slide_data.bullet_points else "")
    if sub:
        _add_text(slide, MARGIN, Inches(3.4), CONTENT_W, Inches(0.9), sub, 20, False, text_color)


def _render_grid(prs, slide_data: SlideContent, theme: SlideTheme, *, cols: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _render_header(slide, slide_data, theme)
    cells = slide_data.grid_cells or []
    n = len(cells)
    rows = (n + cols - 1) // cols
    gap = Inches(0.12)
    grid_top = Inches(1.45)
    grid_h = Inches(5.7)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = (grid_h - gap * (rows - 1)) / rows

    for i, cell in enumerate(cells):
        row, col = divmod(i, cols)
        x = MARGIN + (cell_w + gap) * col
        y = grid_top + (cell_h + gap) * row
        draw_grid_cell(slide, cell, x, y, cell_w, cell_h, theme)


def _render_two_column(prs, slide_data: SlideContent, theme: SlideTheme) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _render_header(slide, slide_data, theme)
    gap = Inches(0.25)
    col_w = (CONTENT_W - gap) / 2
    _draw_bullet_column(slide, MARGIN, Inches(1.45), col_w, Inches(5.7), slide_data.bullet_points, theme)
    if slide_data.visual:
        _draw_visual(slide, slide_data.visual, theme, MARGIN + col_w + gap, Inches(1.45), col_w, Inches(5.7))


def _draw_bullet_column(slide, left, top, width, height, bullets: list[str], theme: SlideTheme) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    row_h = _emu_inches(height) / max(len(bullets), 1)
    for i, bullet in enumerate(bullets[:5]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.alignment = PP_ALIGN.LEFT
        pt = calculate_font_size(bullet, _emu_inches(width), row_h, 15)
        font = p.font
        font.name = FONT_BODY
        font.size = Pt(pt)
        font.color.rgb = hex_to_rgb(theme.dominant_color)


def _draw_visual(slide, visual: VisualElement, theme: SlideTheme, left, top, width, height) -> None:
    dispatch = {
        "BIG_STAT": _draw_big_stat,
        "STEP_FLOW": _draw_step_flow,
        "MATURITY_BAR": _draw_maturity_bar,
        "ICON_LIST": _draw_icon_list,
        "COMPARISON": _draw_comparison,
    }
    dispatch[visual.kind](slide, visual, theme, left, top, width, height)


def _draw_big_stat(slide, visual, theme, left, top, width, height) -> None:
    _add_rect(slide, left, top, width, height, theme.support_color)
    stat_color = _pick_readable_text(theme.support_color, theme.accent_color, theme.dominant_color)
    _add_text(slide, left + Inches(0.15), top + Inches(1.2), width - Inches(0.3), Inches(2.0),
              visual.stat_value or "—", 64, True, stat_color, PP_ALIGN.CENTER)
    _add_text(slide, left + Inches(0.15), top + Inches(3.5), width - Inches(0.3), Inches(1.5),
              visual.stat_label or "", 18, False, theme.dominant_color, PP_ALIGN.CENTER)


def _draw_step_flow(slide, visual, theme, left, top, width, height) -> None:
    items = visual.items or []
    descs = visual.item_descriptions or [""] * len(items)
    n = max(len(items), 1)
    gap = Inches(0.1)
    step_h = (height - gap * (n - 1)) / n
    icon_w = Inches(0.75)
    for i, (item, desc) in enumerate(zip(items[:5], descs[:5])):
        y = top + (step_h + gap) * i
        _add_rect(slide, left, y, width, step_h, theme.support_color)
        emoji, label = _extract_step_emoji(item, i)
        _add_emoji_icon(
            slide, left + Inches(0.08), y, icon_w, step_h,
            emoji, color_hex=theme.accent_color,
        )
        text_x = left + icon_w + Inches(0.05)
        text_w = width - icon_w - Inches(0.18)
        body = f"{i + 1:02d}  {label}" if label else f"{i + 1:02d}"
        if desc:
            body = f"{body}\n{desc}"
        _add_text(
            slide, text_x, y + Inches(0.08), text_w, step_h - Inches(0.12),
            body, 14, False, theme.dominant_color,
        )


def _draw_maturity_bar(slide, visual, theme, left, top, width, height) -> None:
    stages = visual.stages or []
    descs = visual.stage_descriptions or [""] * len(stages)
    n = max(len(stages), 1)
    bar_h = Inches(1.0)
    bar_top = top + Inches(0.8)
    seg_w = width / n
    for i, (stage, desc) in enumerate(zip(stages[:5], descs[:5])):
        ratio = i / max(n - 1, 1)
        color = _blend_hex(theme.support_color, theme.dominant_color, ratio)
        x = left + seg_w * i
        _add_rect(slide, x, bar_top, seg_w - Inches(0.04), bar_h, color, rounded=False)
        tc = "#FFFFFF" if ratio > 0.45 else theme.dominant_color
        _add_text(slide, x + Inches(0.04), bar_top + Inches(0.1), seg_w - Inches(0.1), Inches(0.4),
                  stage, 13, True, tc, PP_ALIGN.CENTER)
        if desc:
            _add_text(slide, x + Inches(0.04), bar_top + Inches(1.15), seg_w - Inches(0.08), Inches(0.8),
                      desc, 11, False, theme.dominant_color)


def _draw_icon_list(slide, visual, theme, left, top, width, height) -> None:
    icons = visual.icon_items or []
    n = max(len(icons), 1)
    gap = Inches(0.1)
    row_h = (height - gap * (n - 1)) / n
    icon_w = Inches(0.85)
    for i, item in enumerate(icons[:5]):
        y = top + (row_h + gap) * i
        row_bg = theme.support_color if i % 2 == 0 else "#FFFFFF"
        _add_rect(slide, left, y, width, row_h, row_bg)
        tc = theme.dominant_color
        _add_emoji_icon(
            slide, left + Inches(0.1), y, icon_w, row_h,
            item.icon, color_hex=theme.accent_color,
        )
        text_x = left + icon_w + Inches(0.08)
        text_w = width - icon_w - Inches(0.2)
        text = f"{item.header}\n{item.body}"
        _add_text(
            slide, text_x, y + Inches(0.06), text_w, row_h - Inches(0.1),
            text, 14, False, tc,
        )


def _draw_comparison(slide, visual, theme, left, top, width, height) -> None:
    card_h = (height - Inches(0.7)) / 2
    _add_rect(slide, left, top, width, card_h, theme.support_color)
    _add_text(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), card_h - Inches(0.2),
              f"Before: {visual.before_label}\n{visual.before_body or ''}", 14, False, theme.dominant_color)
    mid = top + card_h + Inches(0.1)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, left + width / 2 - Inches(0.25), mid, Inches(0.5), Inches(0.45))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = hex_to_rgb(theme.accent_color)
    arrow.line.fill.background()
    after_top = mid + Inches(0.5)
    _add_rect(slide, left, after_top, width, card_h, theme.dominant_color)
    _add_text(slide, left + Inches(0.15), after_top + Inches(0.12), width - Inches(0.3), card_h - Inches(0.2),
              f"After: {visual.after_label}\n{visual.after_body or ''}", 14, True, "#FFFFFF")


def _render_full_step_flow(prs, slide_data: SlideContent, theme: SlideTheme) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _render_header(slide, slide_data, theme)
    if slide_data.visual:
        _draw_step_flow(slide, slide_data.visual, theme, MARGIN, Inches(1.45), CONTENT_W, Inches(5.7))


def _render_maturity_full(prs, slide_data: SlideContent, theme: SlideTheme) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _render_header(slide, slide_data, theme)
    if slide_data.visual:
        _draw_maturity_bar(slide, slide_data.visual, theme, MARGIN, Inches(1.45), CONTENT_W, Inches(5.7))
