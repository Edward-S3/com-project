"""
office_files.py — Excel / Word / PowerPoint の読み込み・編集・出力
"""
from __future__ import annotations

import base64
import io
import json
import re
from datetime import datetime, timedelta, timezone

import pandas as pd

JST = timezone(timedelta(hours=9))

OFFICE_EXTENSIONS = frozenset({"xlsx", "docx", "pptx"})
OFFICE_MIME_MAP = {
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

DOCX_OUTPUT_RE = re.compile(
    r"(docx|\.docx|word形式|ワード形式|Word形式|Wordで|docx形式|docxで|Wordファイル)",
    re.IGNORECASE,
)
PPTX_OUTPUT_RE = re.compile(
    r"(pptx|\.pptx|powerpoint|パワーポイント|PowerPoint|スライド形式|"
    r"プレゼン資料|プレゼンテーション|pptx形式|pptxで|スライドを作成|スライド生成)",
    re.IGNORECASE,
)

_OFFICE_MARKER_RE = re.compile(r"\n?<!--NAI_OFFICE:([A-Za-z0-9+/=]+)-->\s*$")
_SLIDE_HEADING_RE = re.compile(r"^#{1,3}\s+(.+)$", re.MULTILINE)


def file_ext(name: str) -> str:
    return name.rsplit(".", 1)[-1].lower() if "." in name else ""


def is_xlsx(name: str, mime: str = "") -> bool:
    return file_ext(name) == "xlsx" or "spreadsheet" in (mime or "")


def is_docx(name: str, mime: str = "") -> bool:
    if file_ext(name) == "docx":
        return True
    mime_l = (mime or "").lower()
    return (
        "wordprocessing" in mime_l
        or mime_l in (
            "application/msword",
            "application/vnd.ms-word",
        )
    )


def is_pptx(name: str, mime: str = "") -> bool:
    return file_ext(name) == "pptx" or "presentation" in (mime or "")


def is_office_file(name: str, mime: str = "") -> bool:
    return is_xlsx(name, mime) or is_docx(name, mime) or is_pptx(name, mime)


def detect_output_format(prompt: str, system_prompt: str = "") -> str | None:
    """プロンプトから docx / pptx 出力意図を検出（pptx を docx より優先）"""
    combined = f"{prompt}\n{system_prompt}"
    if PPTX_OUTPUT_RE.search(combined):
        return "pptx"
    if DOCX_OUTPUT_RE.search(combined):
        return "docx"
    return None


def output_filename(fmt: str) -> str:
    ts = datetime.now(JST).strftime("%Y%m%d_%H%M%S")
    return f"generated_{ts}.{fmt}"


# ── Excel ────────────────────────────────────────────────

def load_xlsx_sheets(data: bytes) -> dict[str, pd.DataFrame]:
    """xlsx をシート名 → DataFrame の辞書に読み込む"""
    book = pd.read_excel(io.BytesIO(data), sheet_name=None, engine="openpyxl")
    return {
        name: df.fillna("")
        for name, df in book.items()
    }


def sheets_to_text(sheets: dict[str, pd.DataFrame], max_rows: int = 500) -> str:
    """LLM 送信用の表形式テキスト"""
    if not sheets:
        return "[Excel: シートがありません]"
    parts: list[str] = []
    for name, df in sheets.items():
        view = df.head(max_rows)
        parts.append(f"### シート: {name}（{len(df)} 行）\n{view.to_csv(index=False)}")
        if len(df) > max_rows:
            parts.append(f"（先頭 {max_rows} 行のみ表示）")
    return "\n\n".join(parts)


def workbook_bytes_from_sheets(sheets: dict[str, pd.DataFrame]) -> bytes:
    """編集済みシートを xlsx バイナリへ"""
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        for name, df in sheets.items():
            safe_name = str(name)[:31] or "Sheet1"
            df.to_excel(writer, sheet_name=safe_name, index=False)
    return buf.getvalue()


def extract_xlsx_text(data: bytes) -> str:
    return sheets_to_text(load_xlsx_sheets(data))


AGG_LABELS: dict[str, str] = {
    "sum": "合計",
    "mean": "平均",
    "count": "件数",
    "min": "最小",
    "max": "最大",
    "nunique": "ユニーク数",
}


def _numeric_series(df: pd.DataFrame, col: str) -> pd.Series:
    """列を数値として解釈（空文字は NaN）"""
    return pd.to_numeric(df[col].replace("", pd.NA), errors="coerce")


def numeric_columns(df: pd.DataFrame) -> list[str]:
    """集計可能な数値列名の一覧"""
    found: list[str] = []
    for col in df.columns:
        s = _numeric_series(df, col)
        if s.notna().any():
            found.append(str(col))
    return found


def sheet_basic_summary(df: pd.DataFrame) -> pd.DataFrame:
    """1シートの数値列について基本統計を返す"""
    rows: list[dict] = []
    for col in numeric_columns(df):
        s = _numeric_series(df, col)
        rows.append({
            "列名": col,
            "件数": int(s.count()),
            "合計": round(float(s.sum()), 4) if s.count() else None,
            "平均": round(float(s.mean()), 4) if s.count() else None,
            "最小": round(float(s.min()), 4) if s.count() else None,
            "最大": round(float(s.max()), 4) if s.count() else None,
        })
    if not rows:
        return pd.DataFrame([{
            "列名": f"（数値列なし / データ行数 {len(df)}）",
            "件数": len(df),
            "合計": None,
            "平均": None,
            "最小": None,
            "最大": None,
        }])
    return pd.DataFrame(rows)


def workbook_overview(sheets: dict[str, pd.DataFrame]) -> pd.DataFrame:
    """全シートの行数・列数・数値列数の一覧"""
    rows = []
    for name, df in sheets.items():
        rows.append({
            "シート名": name,
            "行数": len(df),
            "列数": len(df.columns),
            "数値列数": len(numeric_columns(df)),
        })
    return pd.DataFrame(rows)


def column_aggregate(df: pd.DataFrame, col: str, func: str) -> pd.DataFrame:
    """1列をシート全体で集計（グループなし）"""
    if col not in df.columns:
        raise ValueError(f"列「{col}」が見つかりません。")
    label = AGG_LABELS.get(func, func)

    if func == "count":
        return pd.DataFrame([{"項目": col, label: len(df)}])
    if func == "nunique":
        return pd.DataFrame([{"項目": col, label: int(df[col].nunique())}])

    s = _numeric_series(df, col)
    if not s.notna().any():
        raise ValueError(f"「{col}」は数値として集計できません。")

    if func == "sum":
        val = float(s.sum())
    elif func == "mean":
        val = float(s.mean())
    elif func == "min":
        val = float(s.min())
    elif func == "max":
        val = float(s.max())
    else:
        raise ValueError(f"未対応の集計: {func}")

    return pd.DataFrame([{"項目": col, label: round(val, 4)}])


def group_aggregate(
    df: pd.DataFrame,
    group_col: str,
    value_col: str | None,
    func: str,
) -> pd.DataFrame:
    """グループ列で集計（value_col 省略時は件数）"""
    if group_col not in df.columns:
        raise ValueError(f"グループ列「{group_col}」が見つかりません。")
    grouped = df.groupby(group_col, dropna=False)
    label = AGG_LABELS.get(func, func)

    if func == "count" or not value_col:
        result = grouped.size().reset_index(name=label)
        return result

    if value_col not in df.columns:
        raise ValueError(f"集計列「{value_col}」が見つかりません。")

    s = _numeric_series(df, value_col)
    work = df.copy()
    work["__agg_val__"] = s
    g = work.groupby(group_col, dropna=False)["__agg_val__"]

    if func == "sum":
        result = g.sum().reset_index(name=label)
    elif func == "mean":
        result = g.mean().reset_index(name=label)
    elif func == "min":
        result = g.min().reset_index(name=label)
    elif func == "max":
        result = g.max().reset_index(name=label)
    elif func == "nunique":
        result = g.nunique().reset_index(name=label)
    else:
        raise ValueError(f"未対応の集計: {func}")

    if label != value_col:
        result = result.rename(columns={label: f"{value_col}_{label}"})
    return result


def execute_agg_instruction(
    sheets: dict[str, pd.DataFrame],
    instruction: dict,
) -> tuple[str, pd.DataFrame]:
    """登録済み集計指示を1件実行し (タイトル, 結果DataFrame) を返す"""
    sheet = instruction.get("sheet", "")
    if sheet not in sheets:
        raise ValueError(f"シート「{sheet}」がありません。")
    df = sheets[sheet]
    func = instruction.get("func", "sum")
    value_col = instruction.get("value_col") or None
    group_col = instruction.get("group_col") or None
    custom = (instruction.get("label") or "").strip()

    if group_col:
        result = group_aggregate(df, group_col, value_col, func)
        auto_title = (
            f"シート「{sheet}」/ {group_col} × {value_col or '（件数）'}"
            f"（{AGG_LABELS.get(func, func)}）"
        )
    elif value_col:
        result = column_aggregate(df, value_col, func)
        auto_title = (
            f"シート「{sheet}」/ {value_col}"
            f"（{AGG_LABELS.get(func, func)}）"
        )
    elif func == "count":
        result = pd.DataFrame([{"シート": sheet, "件数": len(df)}])
        auto_title = f"シート「{sheet}」/ 行数"
    else:
        raise ValueError("集計対象の項目を指定してください。")

    return custom or auto_title, result


def execute_agg_instructions(
    sheets: dict[str, pd.DataFrame],
    instructions: list[dict],
) -> list[tuple[str, pd.DataFrame]]:
    """複数の集計指示を順に実行"""
    out: list[tuple[str, pd.DataFrame]] = []
    for inst in instructions:
        try:
            out.append(execute_agg_instruction(sheets, inst))
        except ValueError as e:
            title = (inst.get("label") or inst.get("sheet", "集計")).strip()
            out.append((f"{title}（エラー）", pd.DataFrame([{"エラー": str(e)}])))
    return out


def dataframe_to_summary_text(df: pd.DataFrame, title: str) -> str:
    if df is None or df.empty:
        return f"### {title}\n（データなし）"
    return f"### {title}\n{df.to_csv(index=False)}"


def sheets_aggregation_text(
    sheets: dict[str, pd.DataFrame],
    *,
    active_sheet: str | None = None,
    group_agg: dict | None = None,
    agg_instructions: list[dict] | None = None,
    include_workbook_overview: bool = True,
    include_sheet_summary: bool = True,
) -> str:
    """LLM 送信用 — 集計結果テキスト"""
    parts: list[str] = ["--- Excel 集計サマリー ---"]
    if include_workbook_overview:
        parts.append(dataframe_to_summary_text(workbook_overview(sheets), "ブック全体"))
    if include_sheet_summary:
        targets = (
            [active_sheet] if active_sheet and active_sheet in sheets
            else list(sheets.keys())
        )
        for name in targets:
            parts.append(
                dataframe_to_summary_text(
                    sheet_basic_summary(sheets[name]),
                    f"シート「{name}」数値列サマリー",
                )
            )
    if agg_instructions:
        for title, result_df in execute_agg_instructions(sheets, agg_instructions):
            parts.append(dataframe_to_summary_text(result_df, title))
    elif group_agg and group_agg.get("result") is not None:
        gsheet = group_agg.get("sheet", "")
        gcol = group_agg.get("group_col", "")
        vcol = group_agg.get("value_col") or "（件数）"
        func = group_agg.get("func", "")
        title = (
            f"グループ集計: シート「{gsheet}」"
            f" / {gcol} × {vcol}（{AGG_LABELS.get(func, func)}）"
        )
        parts.append(dataframe_to_summary_text(group_agg["result"], title))
    return "\n\n".join(parts)


# ── Word ─────────────────────────────────────────────────

def extract_docx_text(data: bytes) -> str:
    if not data:
        return "[Word: ファイルが空です]"
    try:
        from docx import Document

        doc = Document(io.BytesIO(data))
        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        body = "\n".join(parts).strip()
        return body or "[Word: テキストを抽出できませんでした]"
    except Exception as ex:
        return f"[Word ファイルのテキスト抽出に失敗しました: {type(ex).__name__}]"


def text_to_docx(text: str) -> bytes:
    """Markdown 風テキストを docx に変換"""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    style = doc.styles["Normal"]
    style.font.size = Pt(11)

    for line in (text or "").splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith(("- ", "* ", "• ")):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        else:
            doc.add_paragraph(stripped)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PowerPoint ───────────────────────────────────────────

def extract_pptx_text(data: bytes) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                parts.append(f"### スライド {i}\n" + "\n".join(texts))
        body = "\n\n".join(parts).strip()
        return body or "[PowerPoint: テキストを抽出できませんでした]"
    except Exception:
        return "[PowerPoint ファイルのテキスト抽出に失敗しました]"


def _parse_slides_from_text(text: str) -> list[tuple[str, list[str]]]:
    """LLM 回答からスライド（タイトル, 箇条書き）を抽出"""
    text = (text or "").strip()
    if not text:
        return [("スライド 1", ["（内容なし）"])]

    slides: list[tuple[str, list[str]]] = []
    current_title = "スライド 1"
    current_bullets: list[str] = []

    def flush() -> None:
        nonlocal current_title, current_bullets
        bullets = current_bullets or [text[:500] if not slides else "（内容なし）"]
        slides.append((current_title, bullets))
        current_bullets = []

    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        heading = _SLIDE_HEADING_RE.match(stripped)
        if heading:
            if current_bullets or len(slides) > 0:
                flush()
            current_title = heading.group(1).strip()
            continue
        if stripped in ("---", "***", "___"):
            if current_bullets:
                flush()
                current_title = f"スライド {len(slides) + 1}"
            continue
        if stripped.startswith(("# ", "## ", "### ")):
            if current_bullets or slides:
                flush()
            current_title = stripped.lstrip("#").strip()
            continue
        bullet = stripped
        for prefix in ("- ", "* ", "• ", "・"):
            if bullet.startswith(prefix):
                bullet = bullet[len(prefix):]
                break
        if re.match(r"^\d+[\.\)、]\s*", bullet):
            bullet = re.sub(r"^\d+[\.\)、]\s*", "", bullet)
        current_bullets.append(bullet)

    if current_bullets or not slides:
        flush()

    return slides[:30]


def text_to_pptx(text: str) -> bytes:
    """テキストから簡易プレゼンテーションを生成"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[1]

    slides_data = _parse_slides_from_text(text)
    for idx, (title, bullets) in enumerate(slides_data):
        if idx == 0:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets[:12]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
        else:
            slide = prs.slides.add_slide(blank_layout)
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            tx.text_frame.text = title
            tx.text_frame.paragraphs[0].font.size = Pt(28)
            bx = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.5), Inches(5))
            tf = bx.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets[:12]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def extract_office_text(name: str, mime: str, data: bytes) -> str:
    """Office ファイル種別に応じたテキスト抽出"""
    if is_xlsx(name, mime):
        return extract_xlsx_text(data)
    if is_docx(name, mime):
        return extract_docx_text(data)
    if is_pptx(name, mime):
        return extract_pptx_text(data)
    return f"[未対応の Office 形式: {name}]"


def attachment_context_suffix(name: str, mime: str, data: bytes) -> str:
    body = extract_office_text(name, mime, data)
    label = file_ext(name).upper() or "Office"
    return f"\n\n--- 添付 {label} ({name}) ---\n{body}"


def pptx_output_instruction() -> str:
    return (
        "\n\n【出力形式】ユーザーは PowerPoint（pptx）形式での成果物を求めています。"
        "回答は各スライドを「## スライドタイトル」で区切り、"
        "本文は箇条書き（- で始める）で記述してください。"
    )


# ── メッセージへの出力ファイル埋め込み ───────────────────

def serialize_office_outputs(outputs: list[dict]) -> str:
    if not outputs:
        return ""
    payload = [
        {
            "name": o["name"],
            "mime": o["mime"],
            "b64": base64.standard_b64encode(o["data"]).decode(),
        }
        for o in outputs
    ]
    encoded = base64.standard_b64encode(
        json.dumps(payload, ensure_ascii=False).encode()
    ).decode()
    return f"\n<!--NAI_OFFICE:{encoded}-->"


def split_content_and_outputs(content: str) -> tuple[str, list[dict]]:
    m = _OFFICE_MARKER_RE.search(content or "")
    if not m:
        return content or "", []
    text = (content or "")[: m.start()].rstrip()
    try:
        payload = json.loads(base64.standard_b64decode(m.group(1)).decode())
        outputs = [
            {
                "name": item["name"],
                "mime": item["mime"],
                "data": base64.standard_b64decode(item["b64"]),
            }
            for item in payload
        ]
        return text, outputs
    except Exception:
        return content or "", []


def build_output_files(fmt: str, text: str) -> list[dict]:
    if fmt == "docx":
        return [{
            "name": output_filename("docx"),
            "mime": OFFICE_MIME_MAP["docx"],
            "data": text_to_docx(text),
        }]
    if fmt == "pptx":
        return [{
            "name": output_filename("pptx"),
            "mime": OFFICE_MIME_MAP["pptx"],
            "data": text_to_pptx(text),
        }]
    return []
