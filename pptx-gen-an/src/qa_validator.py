"""V2: 生成 PPTX の品質検証（テキストオーバーフロー検出）"""
from __future__ import annotations

import io
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Pt


@dataclass
class QAIssue:
    slide_index: int
    message: str


@dataclass
class QAResult:
    ok: bool
    issues: list[QAIssue] = field(default_factory=list)

    @property
    def warnings(self) -> list[str]:
        return [f"スライド {i.slide_index}: {i.message}" for i in self.issues]


# おおよその1行あたり文字数（16pt・幅12インチ想定）
CHARS_PER_LINE = 42
LINE_HEIGHT_PT = 20


def _estimate_lines(text: str, font_size_pt: float) -> int:
    if not text:
        return 0
    scale = 16.0 / max(font_size_pt, 10.0)
    chars_per_line = int(CHARS_PER_LINE * scale)
    lines = 0
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        lines += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
    return lines


def validate_pptx_bytes(data: bytes) -> QAResult:
    """テキストボックスの高さと文字数の整合性を検証"""
    issues: list[QAIssue] = []
    prs = Presentation(io.BytesIO(data))

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full_text = "\n".join(p.text for p in tf.paragraphs if p.text)
            if not full_text.strip():
                continue

            font_size = 16.0
            for p in tf.paragraphs:
                if p.font.size:
                    font_size = float(p.font.size.pt)
                    break

            box_height_pt = float(shape.height.pt)
            max_lines = max(1, int(box_height_pt / LINE_HEIGHT_PT))
            est_lines = _estimate_lines(full_text, font_size)

            if est_lines > max_lines + 1:
                issues.append(
                    QAIssue(
                        slide_index=slide_idx,
                        message=(
                            f"テキストがはみ出す可能性があります "
                            f"（推定{est_lines}行 / 枠{max_lines}行, {len(full_text)}文字）"
                        ),
                    )
                )

    return QAResult(ok=len(issues) == 0, issues=issues)


def try_reduce_overflow(data: bytes, warnings: list[str]) -> bytes:
    """
    オーバーフロー警告がある場合、該当テキストのフォントを1段階下げて再保存。
    修正できない場合は元データを返す。
    """
    if not warnings:
        return data

    prs = Presentation(io.BytesIO(data))
    changed = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for p in tf.paragraphs:
                if not p.text or not p.font.size:
                    continue
                if p.font.size.pt > 14:
                    p.font.size = Pt(max(14, int(p.font.size.pt) - 2))
                    changed = True

    if not changed:
        return data
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
