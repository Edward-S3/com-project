import streamlit as st
import os
import io
import json
import requests
from dotenv import load_dotenv

# 環境変数の読み込み
load_dotenv()

import tempfile
import google.generativeai as genai
import pandas as pd
import docx
from docx.shared import Inches
import pptx
import pdfplumber
from pdf2image import convert_from_path
import pytesseract
from PIL import Image

# グラフ描画用
import matplotlib.pyplot as plt
import numpy as np
import japanize_matplotlib  # 日本語文字化け防止
import openpyxl

# PDF生成用ライブラリ
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.lib.utils import ImageReader

# --- セッションステートの初期化 ---
for key in ['evaluation_done', 'evaluation_result', 'chat_history', 'discussion_display', 'employee_info', 'model_answer', 'model_answer_format']:
    if key not in st.session_state:
        st.session_state[key] = False if key == 'evaluation_done' else None if 'format' in key or 'result' in key or 'answer' in key else [] if 'history' in key or 'display' in key else {}

def reset_app():
    for key in list(st.session_state.keys()):
        del st.session_state[key]
    st.rerun()

# --- バックエンド: テキスト抽出ロジック ---
def extract_text_from_file(uploaded_file):
    filename = uploaded_file.name
    ext = filename.split('.')[-1].lower()
    text = f"【ファイル名: {filename}】\n"
    try:
        if ext == 'txt':
            text += uploaded_file.getvalue().decode('utf-8')
        elif ext in ['csv', 'xlsx']:
            df = pd.read_csv(uploaded_file) if ext == 'csv' else pd.read_excel(uploaded_file)
            text += df.to_string()
        elif ext == 'docx':
            doc = docx.Document(uploaded_file)
            text += "\n".join([p.text for p in doc.paragraphs])
        elif ext == 'pptx':
            ppt = pptx.Presentation(uploaded_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == 'pdf':
            pdf_extracted_text = ""
            with pdfplumber.open(uploaded_file) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pdf_extracted_text += page_text + "\n"
            if len(pdf_extracted_text.strip()) < 200:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
                    tmp.write(uploaded_file.getvalue())
                    tmp_path = tmp.name
                images = convert_from_path(tmp_path)
                for img in images:
                    text += pytesseract.image_to_string(img, lang='jpn')
                os.remove(tmp_path)
            else:
                text += pdf_extracted_text
    except Exception as e:
        text += f"\n[ファイルの読み取り中にエラーが発生しました: {str(e)}]\n"
    return text

# --- バックエンド: レーダーチャート生成ロジック ---
def create_radar_chart(radar_scores):
    if not radar_scores:
        radar_scores = {"論理・説得力": 3, "専門知識": 3, "実務適合性": 3, "課題解決力": 3, "構成・表現力": 3}
    
    labels = list(radar_scores.keys())
    stats = list(radar_scores.values())
    
    angles = np.linspace(0, 2 * np.pi, len(labels), endpoint=False).tolist()
    stats = np.concatenate((stats, [stats[0]]))
    angles = np.concatenate((angles, [angles[0]]))
    
    fig, ax = plt.subplots(figsize=(5, 5), subplot_kw=dict(polar=True))
    ax.fill(angles, stats, color='skyblue', alpha=0.4)
    ax.plot(angles, stats, color='#1f77b4', linewidth=2)
    ax.set_yticklabels([])
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, fontsize=12)
    ax.set_ylim(0, 5)
    
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=False, facecolor='white')
    buf.seek(0)
    plt.close(fig)
    return buf

# --- バックエンド: 評価結果のエクスポート用ロジック ---
def export_eval_docx(emp_info, result, standard):
    doc = docx.Document()
    doc.add_heading(f"評価結果: {emp_info.get('氏名', '')} 様", level=0)
    doc.add_paragraph(f"適用基準: {standard}")
    
    # チャット画像を挿入
    chart_buf = create_radar_chart(result.get("radar_scores", {}))
    doc.add_picture(chart_buf, width=Inches(4.0))
    
    doc.add_heading("【被評価者情報】", level=2)
    for k, v in emp_info.items(): doc.add_paragraph(f"{k}: {v}")
    doc.add_heading("【評価サマリー】", level=2)
    doc.add_paragraph(f"総合評価点: {result.get('total_score', 0)} 点")
    doc.add_heading("【評価の詳細】", level=2)
    doc.add_paragraph(f"主要トピック: {result.get('topic', '')}")
    doc.add_heading("■ Good Point", level=3)
    doc.add_paragraph(result.get('good_points', ''))
    doc.add_heading("■ Bad Point", level=3)
    doc.add_paragraph(result.get('bad_points', ''))
    doc.add_heading("【評価根拠】", level=2)
    doc.add_paragraph(result.get('evaluation_reason', ''))
    doc.add_heading("【改善案】", level=2)
    doc.add_paragraph(result.get('improvement_plan', ''))
    
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

def export_eval_xlsx(emp_info, result, standard):
    data = {
        "氏名": emp_info.get("氏名"), "適用基準": standard, "総合評価点": result.get("total_score"),
        "主要トピック": result.get("topic"), "Good Point": result.get("good_points"),
        "Bad Point": result.get("bad_points"), "評価根拠": result.get("evaluation_reason"), "改善案": result.get("improvement_plan")
    }
    df = pd.DataFrame([data])
    bio = io.BytesIO()
    
    with pd.ExcelWriter(bio, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='評価結果')
        worksheet = writer.sheets['評価結果']
        
        # セルの幅を少し広げる
        for col in worksheet.columns:
            worksheet.column_dimensions[col[0].column_letter].width = 20
            
        chart_buf = create_radar_chart(result.get("radar_scores", {}))
        img = openpyxl.drawing.image.Image(chart_buf)
        # データの真下に画像を配置
        worksheet.add_image(img, 'A4')
        
    return bio.getvalue()

def export_eval_pdf(emp_info, result, standard):
    bio = io.BytesIO()
    c = canvas.Canvas(bio, pagesize=A4)
    pdfmetrics.registerFont(UnicodeCIDFont('HeiseiKakuGo-W5'))
    
    width, height = A4
    y = height - 50
    margin = 50
    font_name = 'HeiseiKakuGo-W5'

    def draw_text(text, font_size=11):
        nonlocal y
        c.setFont(font_name, font_size)
        for paragraph in str(text).split('\n'):
            line = ""
            for char in paragraph:
                if pdfmetrics.stringWidth(line + char, font_name, font_size) <= (width - 2 * margin):
                    line += char
                else:
                    if y < 50:
                        c.showPage()
                        c.setFont(font_name, font_size)
                        y = height - 50
                    c.drawString(margin, y, line)
                    y -= 16
                    line = char
            if line:
                if y < 50:
                    c.showPage()
                    c.setFont(font_name, font_size)
                    y = height - 50
                c.drawString(margin, y, line)
                y -= 16

    draw_text(f"評価結果: {emp_info.get('氏名', '')} 様", 16)
    y -= 5
    draw_text(f"適用基準: {standard}", 10)
    y -= 15
    
    # チャット画像の挿入
    chart_buf = create_radar_chart(result.get("radar_scores", {}))
    # 画像のスペースが足りない場合は改ページ
    if y < 250:
        c.showPage()
        y = height - 50
    
    # 画像を描画 (中央付近に配置)
    c.drawImage(ImageReader(chart_buf), margin + 80, y - 250, width=250, height=250)
    y -= 270 # 画像の分だけY座標を下げる

    draw_text("【評価サマリー】", 12)
    draw_text(f"  総合評価点: {result.get('total_score', 0)} 点")
    y -= 15
    draw_text("【評価の詳細】", 12)
    draw_text(f"  主要トピック: {result.get('topic', '')}")
    y -= 10
    draw_text("■ Good Point", 11)
    draw_text(result.get('good_points', ''))
    y -= 10
    draw_text("■ Bad Point", 11)
    draw_text(result.get('bad_points', ''))
    y -= 15
    draw_text("【評価根拠】", 12)
    draw_text(result.get('evaluation_reason', ''))
    y -= 15
    draw_text("【改善案】", 12)
    draw_text(result.get('improvement_plan', ''))

    c.save()
    return bio.getvalue()

# --- バックエンド: 模範解答（Word）生成ロジック ---
def create_docx_from_text(text):
    doc = docx.Document()
    doc.add_heading('100点満点 模範解答', level=0)
    for line in text.split('\n'):
        if line.startswith('# '): doc.add_heading(line[2:], level=1)
        elif line.startswith('## '): doc.add_heading(line[3:], level=2)
        elif line.startswith('### '): doc.add_heading(line[4:], level=3)
        elif line.strip() == '': continue
        else: doc.add_paragraph(line)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# --- バックエンド: LLMバックエンド（ローカルLLM / API）抽象化 ---
LLM_PROVIDER_LOCAL = "ローカルLLM (Ollama)"
LLM_PROVIDER_GEMINI = "API (Gemini)"

DEFAULT_OLLAMA_URL = os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434")
DEFAULT_OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "llama3.2:latest")
GEMINI_MODEL_CHOICES = ["gemini-3.5-flash", "gemini-2.5-flash", "gemini-2.5-pro", "gemini-1.5-flash", "gemini-1.5-pro"]

def init_llm_config():
    """LLM接続設定をセッションステートに初期化する。"""
    if 'llm_config' not in st.session_state or not isinstance(st.session_state.llm_config, dict):
        env_provider = os.environ.get("LLM_PROVIDER", "local").lower()
        provider = LLM_PROVIDER_LOCAL if env_provider == "local" else LLM_PROVIDER_GEMINI
        st.session_state.llm_config = {
            "provider": provider,
            "ollama_url": DEFAULT_OLLAMA_URL,
            "ollama_model": DEFAULT_OLLAMA_MODEL,
            "gemini_api_key": os.environ.get("GEMINI_API_KEY", ""),
            "gemini_model": GEMINI_MODEL_CHOICES[0],
        }

def list_ollama_models(base_url):
    """Ollamaサーバーから利用可能なモデル名一覧を取得する。"""
    try:
        r = requests.get(f"{base_url.rstrip('/')}/api/tags", timeout=3)
        r.raise_for_status()
        return [m.get("name") for m in r.json().get("models", []) if m.get("name")]
    except Exception:
        return []

def llm_is_ready(cfg=None):
    """選択中のバックエンドが利用可能かを判定する。"""
    cfg = cfg or st.session_state.get('llm_config', {})
    if cfg.get("provider") == LLM_PROVIDER_LOCAL:
        return bool(list_ollama_models(cfg.get("ollama_url", "")))
    return bool(cfg.get("gemini_api_key"))

def _ollama_chat(base_url, model, messages, system_prompt=None, temperature=0.2, json_mode=False):
    """Ollama /api/chat を呼び出してテキストを返す。"""
    payload_messages = []
    if system_prompt:
        payload_messages.append({"role": "system", "content": system_prompt})
    for m in messages:
        role = "assistant" if m["role"] in ("assistant", "model") else m["role"]
        payload_messages.append({"role": role, "content": m["content"]})

    payload = {
        "model": model,
        "messages": payload_messages,
        "stream": False,
        "options": {"temperature": float(temperature)},
    }
    if json_mode:
        payload["format"] = "json"

    r = requests.post(f"{base_url.rstrip('/')}/api/chat", json=payload, timeout=600)
    r.raise_for_status()
    data = r.json()
    return data.get("message", {}).get("content", "")

def _gemini_chat(api_key, model_name, messages, system_prompt=None, temperature=0.2, json_mode=False):
    """Gemini APIを呼び出してテキストを返す。"""
    if not api_key:
        raise RuntimeError("Gemini APIキーが設定されていません。")
    genai.configure(api_key=api_key)

    gemini_history = []
    for m in messages:
        role = "user" if m["role"] == "user" else "model"
        gemini_history.append({"role": role, "parts": [m["content"]]})

    gen_config = {"temperature": float(temperature)}
    if json_mode:
        gen_config["response_mime_type"] = "application/json"

    model = genai.GenerativeModel(
        model_name=model_name,
        system_instruction=system_prompt,
        generation_config=gen_config,
    )
    response = model.generate_content(gemini_history)
    return response.text

def llm_chat(messages, system_prompt=None, temperature=0.2, json_mode=False, cfg=None):
    """設定に応じてローカルLLM / API を切り替えて応答テキストを返す。

    messages は [{"role": "user"|"assistant", "content": str}, ...] 形式。
    """
    cfg = cfg or st.session_state.get('llm_config', {})
    if cfg.get("provider") == LLM_PROVIDER_LOCAL:
        return _ollama_chat(
            cfg.get("ollama_url", DEFAULT_OLLAMA_URL),
            cfg.get("ollama_model", DEFAULT_OLLAMA_MODEL),
            messages, system_prompt=system_prompt,
            temperature=temperature, json_mode=json_mode,
        )
    return _gemini_chat(
        cfg.get("gemini_api_key", ""),
        cfg.get("gemini_model", GEMINI_MODEL_CHOICES[0]),
        messages, system_prompt=system_prompt,
        temperature=temperature, json_mode=json_mode,
    )

def parse_json_response(text):
    """LLM応答からJSONを抽出してdictに変換する（ローカルLLMの揺らぎ対策込み）。"""
    if text is None:
        raise ValueError("LLM応答が空です。")
    s = text.strip()
    if s.startswith("```"):
        s = s.strip("`")
        if s.lower().startswith("json"):
            s = s[4:]
        s = s.strip()
    try:
        return json.loads(s)
    except json.JSONDecodeError:
        start = s.find("{")
        end = s.rfind("}")
        if start != -1 and end != -1 and end > start:
            return json.loads(s[start:end + 1])
        raise

# --- バックエンド: システムプロンプト定義 ---
def get_system_prompt(eval_standard, industry):
    industry_instruction = f"4. 【業種特性の考慮】: 所属する『{industry}』の特性を深く考慮して評価を行ってください。\n" if industry else ""
    base_prompt = f"""
【対話と再評価のルール】
反論や補足説明を真摯に吟味し、妥当であれば躊躇なく点数や評価内容を修正（加点）してください。
【評価基準】
1. 基準設定: 年齢や等級に応じた水準を基準とし、優れた工夫があれば加点。
2. 減点: 良い点は積極的に評価し、不備や矛盾はBad Pointとして抽出（1つにつき5〜10点減点）。
{industry_instruction}
【出力要件】
以下のキーを必ず持つJSON形式で出力してください。
{{
  "ai_message": "あなたからのメッセージ。",
  "total_score": 75,
  "raw_scale_1_to_5": 3,
  "radar_scores": {{
    "論理・説得力": 4,
    "専門知識": 3,
    "実務適合性": 4,
    "課題解決力": 3,
    "構成・表現力": 5
  }},
  "topic": "〇〇について",
  "achievement_level": "70%",
  "good_points": "・〇〇\\n・〇〇",
  "bad_points": "・【マイナス10点】〇〇の論理が破綻している",
  "evaluation_reason": "...",
  "improvement_plan": "..."
}}
※ radar_scores は必ず上記5つのキー（完全一致）とし、値は1〜5の整数で評価してください。
"""
    if eval_standard == "トップレベル（厳格・最高水準）":
        return "あなたは厳格かつ極めて論理的な「成果物評価AIアーキテクト」です。\n" + base_prompt.replace("5〜10点減点", "10〜20点減点")
    else:
        return "あなたは社員の育成と適正な評価を支援する「公平かつ建設的な成果物評価AIアシスタント」です。\n" + base_prompt

# --- バックエンド: 模範解答生成ロジック ---
def generate_model_answer(employee_info, eval_result, format_type, user_request=""):
    request_prompt = f"\n【特別な要望】\n{user_request}\n※上記の要望に合わせ、トーン＆マナーを最適化してください。" if user_request else ""
    base_prompt = f"""
    あなたは超一流のコンサルタントです。以下の【被評価者情報】と、AIによる【評価結果】を踏まえ、
    すべてのBad Pointを完全に解消し、改善案を100%反映した「100点満点の模範解答」を作成してください。{request_prompt}
    【被評価者情報】\n{json.dumps(employee_info, ensure_ascii=False)}
    【評価結果】\n{json.dumps(eval_result, ensure_ascii=False)}
    """

    if "GASコード" in format_type:
        format_prompt = "【出力形式】Googleスライド自動生成用のGASコードのみを出力。関数名は `function createModelPresentation() { ... }`。テキストと「画像挿入用プレースホルダー枠」を配置し、枠内に画像生成AI用のプロンプト案を記載してください。Markdownブロック不可。"
    elif "Word" in format_type:
        format_prompt = "【出力形式】Wordドキュメント変換用の構造化Markdownテキスト。見出しを活用し、「【推奨画像プロンプト】: 〇〇」という挿絵提案を含めてください。"
    else:
        format_prompt = "【出力形式】実務的なMarkdownテキスト形式。"

    full_prompt = base_prompt + "\n" + format_prompt
    ans = llm_chat(
        messages=[{"role": "user", "content": full_prompt}],
        temperature=0.4,
    )

    if "GASコード" in format_type:
        ans = ans.replace("```javascript", "").replace("```js", "").replace("```", "").strip()
    return ans

# --- フロントエンド: UI実装 ---
st.set_page_config(page_title="社員成果物 自動評価ツール", layout="wide")
st.title("📑 社員成果物 自動評価ツール")

# LLM接続設定の初期化
init_llm_config()

# --- サイドバー: LLMバックエンド切替 ---
with st.sidebar:
    st.header("⚙️ LLM設定")
    cfg = st.session_state.llm_config

    provider = st.radio(
        "使用するLLM",
        [LLM_PROVIDER_LOCAL, LLM_PROVIDER_GEMINI],
        index=0 if cfg["provider"] == LLM_PROVIDER_LOCAL else 1,
        help="ローカルLLM（Ollama）と外部API（Gemini）を切替できます。",
    )
    cfg["provider"] = provider

    if provider == LLM_PROVIDER_LOCAL:
        cfg["ollama_url"] = st.text_input(
            "Ollama ベースURL",
            value=cfg.get("ollama_url", DEFAULT_OLLAMA_URL),
        )
        available_models = list_ollama_models(cfg["ollama_url"])
        if available_models:
            current_model = cfg.get("ollama_model", DEFAULT_OLLAMA_MODEL)
            if current_model not in available_models:
                current_model = available_models[0]
            cfg["ollama_model"] = st.selectbox(
                "モデル",
                available_models,
                index=available_models.index(current_model),
            )
            st.success(f"✅ Ollama接続OK（{len(available_models)}モデル利用可）")
        else:
            cfg["ollama_model"] = st.text_input(
                "モデル名",
                value=cfg.get("ollama_model", DEFAULT_OLLAMA_MODEL),
                help="例: llama3.2:latest / gemma3:4b / qwen2.5:7b など",
            )
            st.error("⚠️ Ollamaサーバーに接続できません。URLとサーバー起動状態を確認してください。")
    else:
        env_key = os.environ.get("GEMINI_API_KEY", "")
        cfg["gemini_api_key"] = env_key
        gem_default = cfg.get("gemini_model", GEMINI_MODEL_CHOICES[0])
        if gem_default not in GEMINI_MODEL_CHOICES:
            gem_default = GEMINI_MODEL_CHOICES[0]
        cfg["gemini_model"] = st.selectbox(
            "Geminiモデル",
            GEMINI_MODEL_CHOICES,
            index=GEMINI_MODEL_CHOICES.index(gem_default),
        )
        if cfg["gemini_api_key"]:
            st.success("✅ Gemini APIキー設定済み（環境変数）")
        else:
            st.warning("⚠️ Gemini APIキーが未設定です。.env に GEMINI_API_KEY を設定してください。")

    st.caption(
        "💡 ローカルLLMは小規模モデルだと評価品質が低下する場合があります。"
        "JSON出力の安定性のため Llama3.1/3.2、qwen2.5、gemma3 などを推奨。"
    )

if not st.session_state.evaluation_done:
    with st.form("evaluation_form"):
        st.subheader("1. 被評価者情報の入力")
        col1, col2, col3 = st.columns(3)
        with col1:
            name = st.text_input("氏名")
            age = st.number_input("年齢", min_value=18, max_value=100, value=30)
            industry = st.text_input("所属企業の業種（例：製造業、コンサル業）", value="")
        with col2:
            gender = st.selectbox("性別", ["男性", "女性", "その他", "回答しない"])
            tenure = st.number_input("勤続年数", min_value=0, max_value=60, value=5)
        with col3:
            department = st.text_input("対象部門", value="総務部")
            grade = st.selectbox("等級", ["等級対象外", "一般社員1-1", "一般社員1-2", "一般社員1-3", "一般社員2", "主任3", "係長4-1", "課長代理4-2", "課長5", "次長6", "部長7", "取締役8"])
            
        st.subheader("2. 評価基準の選択")
        eval_standard = st.radio("AIの評価基準（厳しさ）", ["中小企業水準（標準・育成フォーカス）", "トップレベル（厳格・最高水準）"])
            
        st.subheader("3. 成果物ファイルのアップロード")
        uploaded_files = st.file_uploader("評価対象のファイル", type=["txt", "csv", "pptx", "docx", "xlsx", "pdf"], accept_multiple_files=True)
        
        submitted = st.form_submit_button("評価を実行する")

    if submitted:
        cfg = st.session_state.llm_config
        if cfg["provider"] == LLM_PROVIDER_GEMINI and not cfg.get("gemini_api_key"):
            st.error("Gemini APIキーが設定されていません。サイドバーから入力してください。")
        elif cfg["provider"] == LLM_PROVIDER_LOCAL and not list_ollama_models(cfg["ollama_url"]):
            st.error("Ollamaサーバーに接続できません。サイドバーのURLとサーバー起動状態を確認してください。")
        elif not name or not uploaded_files:
            st.error("氏名とファイルを指定してください。")
        else:
            with st.status("評価プロセスを実行中...", expanded=True) as status:
                try:
                    combined_text = "".join([extract_text_from_file(f) + "\n\n" for f in uploaded_files])
                    employee_info = {"氏名": name, "年齢": age, "性別": gender, "勤続年数": tenure, "対象部門": department, "等級": grade, "業種": industry if industry else "未指定"}

                    system_prompt = get_system_prompt(eval_standard, industry)
                    user_prompt = f"【被評価者情報】\n{json.dumps(employee_info, ensure_ascii=False, indent=2)}\n\n【抽出された成果物テキスト】\n{combined_text}"

                    response_text = llm_chat(
                        messages=[{"role": "user", "content": user_prompt}],
                        system_prompt=system_prompt,
                        temperature=0.2,
                        json_mode=True,
                    )

                    result = parse_json_response(response_text)
                    st.session_state.update({
                        'chat_history': [
                            {"role": "user", "content": user_prompt},
                            {"role": "assistant", "content": response_text},
                        ],
                        'evaluation_result': result, 'employee_info': employee_info,
                        'eval_name': name, 'eval_standard_used': eval_standard, 'industry_used': industry, 'evaluation_done': True,
                        'discussion_display': [{"role": "assistant", "text": result.get("ai_message", "評価が完了しました。")}]
                    })
                    st.rerun()
                except Exception as e:
                    st.error(f"エラー: {str(e)}")

else:
    # --- 結果表示画面 ---
    result = st.session_state.evaluation_result
    name = st.session_state.eval_name
    st.header(f"📊 評価結果: {name} 様")
    st.write(f"適用基準: **{st.session_state.eval_standard_used}**")
    
    col_score, col_chart = st.columns([1, 2])
    with col_score:
        st.metric("総合評価点", f"{result.get('total_score', 0)} 点")
        st.metric("評価指数", f"レベル {result.get('raw_scale_1_to_5', 3) + 2}")
    
    with col_chart:
        # UIにもレーダーチャートを表示
        if "radar_scores" in result:
            chart_img = create_radar_chart(result["radar_scores"])
            st.image(chart_img, use_container_width=False, width=350)
    
    st.subheader("📝 評価の詳細")
    col_g, col_b = st.columns(2)
    with col_g: st.success("**Good Point**\n\n" + result.get('good_points', 'N/A'))
    with col_b: st.error("**Bad Point**\n\n" + result.get('bad_points', 'N/A'))
    
    st.info(f"**🔎 評価根拠:**\n{result.get('evaluation_reason', 'N/A')}")
    st.warning(f"**💡 改善案:**\n{result.get('improvement_plan', 'N/A')}")
    st.markdown("---")
    
    # --- 評価結果のダウンロードセクション ---
    st.subheader("📥 評価結果のダウンロード (グラフ付き)")
    dl_col1, dl_col2, dl_col3 = st.columns(3)
    with dl_col1:
        pdf_data = export_eval_pdf(st.session_state.employee_info, result, st.session_state.eval_standard_used)
        st.download_button("📕 PDF形式 (.pdf)", data=pdf_data, file_name=f"評価結果_{name}様.pdf", mime="application/pdf", use_container_width=True)
    with dl_col2:
        docx_data = export_eval_docx(st.session_state.employee_info, result, st.session_state.eval_standard_used)
        st.download_button("📘 Word形式 (.docx)", data=docx_data, file_name=f"評価結果_{name}様.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
    with dl_col3:
        xlsx_data = export_eval_xlsx(st.session_state.employee_info, result, st.session_state.eval_standard_used)
        st.download_button("📗 Excel形式 (.xlsx)", data=xlsx_data, file_name=f"評価結果_{name}様.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True)

    st.markdown("---")
    
    # --- 議論（チャット）セクション ---
    st.subheader("💬 AIと議論する / 評価に反論する")
    for msg in st.session_state.discussion_display: st.chat_message(msg["role"]).write(msg["text"])
    
    if user_rebuttal := st.chat_input("例：この資料はあくまで社内ドラフト用なので、コスト要件は意図的に省いています。"):
        st.session_state.discussion_display.append({"role": "user", "text": user_rebuttal})
        st.chat_message("user").write(user_rebuttal)
        with st.spinner("AIが反論を吟味し、再評価を行っています..."):
            try:
                rebuttal_prompt = f"【反論・補足】\n{user_rebuttal}\n論理的に吟味し、必要に応じて評価を修正してください。必ずJSON形式で出力。"
                st.session_state.chat_history.append({"role": "user", "content": rebuttal_prompt})
                system_prompt = get_system_prompt(st.session_state.eval_standard_used, st.session_state.industry_used)
                response_text = llm_chat(
                    messages=st.session_state.chat_history,
                    system_prompt=system_prompt,
                    temperature=0.2,
                    json_mode=True,
                )
                new_result = parse_json_response(response_text)
                st.session_state.chat_history.append({"role": "assistant", "content": response_text})
                st.session_state.evaluation_result = new_result
                st.session_state.discussion_display.append({"role": "assistant", "text": new_result.get("ai_message", "評価を更新しました。")})
                st.rerun()
            except Exception as e:
                st.error(f"再評価中にエラーが発生しました: {str(e)}")

    st.markdown("---")
    
    # --- 模範解答生成セクション ---
    st.subheader("💯 100点満点の模範解答を生成する")
    user_request = st.text_area("📝 資料のターゲット層・場面・目的などの要望（任意）", placeholder="例：役員向けにコスト対効果を強調してほしい。")
    output_format = st.radio("出力形式を選択", ["スライド形式（GASコード出力）", "ドキュメント形式（Word .docx出力）", "テキスト形式（Markdown出力）"], horizontal=True)
    
    if st.button("✨ 模範解答を生成する", type="primary"):
        with st.spinner(f"【{output_format}】で模範解答を生成中..."):
            try:
                model_answer = generate_model_answer(st.session_state.employee_info, st.session_state.evaluation_result, output_format, user_request)
                st.session_state.model_answer = model_answer
                st.session_state.model_answer_format = output_format
                st.success("生成完了！以下のボタンからダウンロードしてください。")
            except Exception as e:
                st.error(f"生成中にエラーが発生しました: {str(e)}")

    if st.session_state.model_answer:
        ans = st.session_state.model_answer
        fmt = st.session_state.model_answer_format
        if "GASコード" in fmt:
            st.download_button("📥 GASコード (.js)", data=ans, file_name=f"model_answer_{name}.js", mime="text/javascript")
        elif "Word" in fmt:
            docx_data = create_docx_from_text(ans)
            st.download_button("📥 Wordドキュメント (.docx)", data=docx_data, file_name=f"model_answer_{name}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        else:
            st.download_button("📥 Markdownテキスト (.md)", data=ans, file_name=f"model_answer_{name}.md", mime="text/markdown")

    st.markdown("---")
    if st.button("🔄 全てクリアして別の人の評価を行う"):
        reset_app()