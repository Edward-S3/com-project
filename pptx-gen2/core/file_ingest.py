"""多形式ファイル取り込み・テキスト化。"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

import pandas as pd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

logger = logging.getLogger(__name__)

SCAN_PDF_CHAR_THRESHOLD = 50


@dataclass
class IngestedSource:
    name: str
    source_type: str
    text: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)
    needs_vision_ocr: bool = False
    media_path: Optional[Path] = None
    mime_type: Optional[str] = None


def ingest_file(path: Path) -> IngestedSource:
    suffix = path.suffix.lower()
    name = path.name

    if suffix in {".txt", ".md"}:
        text = path.read_text(encoding="utf-8", errors="replace")
        return IngestedSource(name, "text", text, {"char_count": len(text)})

    if suffix == ".csv":
        df = pd.read_csv(path)
        text = df.to_string(index=False)
        return IngestedSource(name, "csv", text, {"rows": len(df), "char_count": len(text)})

    if suffix == ".pdf":
        reader = PdfReader(str(path))
        pages = [p.extract_text() or "" for p in reader.pages]
        text = "\n\n".join(pages)
        needs_ocr = len(text.strip()) < SCAN_PDF_CHAR_THRESHOLD
        return IngestedSource(
            name,
            "pdf",
            text,
            {"pages": len(reader.pages), "char_count": len(text)},
            needs_vision_ocr=needs_ocr,
        )

    if suffix == ".docx":
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else ""
                parts.append(f"[{style}] {para.text}")
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        text = "\n".join(parts)
        return IngestedSource(name, "docx", text, {"char_count": len(text)})

    if suffix == ".pptx":
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            chunks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
            note = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            slides_text.append(f"--- スライド {i} ---\n" + "\n".join(chunks) + f"\n[ノート] {note}")
        text = "\n\n".join(slides_text)
        return IngestedSource(name, "pptx", text, {"slides": len(prs.slides), "char_count": len(text)})

    if suffix in {".xlsx", ".xlsm"}:
        xl = pd.ExcelFile(path)
        parts = []
        for sheet in xl.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet)
            parts.append(f"=== シート: {sheet} ===\n{df.to_string(index=False)}")
        text = "\n\n".join(parts)
        return IngestedSource(name, "xlsx", text, {"sheets": len(xl.sheet_names), "char_count": len(text)})

    if suffix in {".mp3", ".m4a"}:
        duration = _media_duration(path)
        return IngestedSource(
            name,
            "audio",
            "",
            {"duration_sec": duration},
            media_path=path,
            mime_type="audio/mpeg" if suffix == ".mp3" else "audio/mp4",
        )

    if suffix in {".mp4", ".m4v"}:
        duration = _media_duration(path)
        return IngestedSource(
            name,
            "video",
            "",
            {"duration_sec": duration},
            media_path=path,
            mime_type="video/mp4",
        )

    raise ValueError(f"非対応形式: {suffix}")


def ingest_files(paths: List[Path]) -> List[IngestedSource]:
    results = []
    for p in paths:
        try:
            results.append(ingest_file(p))
        except Exception as exc:
            logger.error("ファイル取り込み失敗 %s: %s", p, exc)
    return results


def _media_duration(path: Path) -> float:
    try:
        from mutagen import File as MutagenFile

        meta = MutagenFile(str(path))
        if meta and meta.info and hasattr(meta.info, "length"):
            return float(meta.info.length)
    except Exception:
        pass
    return 0.0


def local_stats(sources: List[IngestedSource]) -> Dict[str, Any]:
    total_chars = sum(len(s.text) for s in sources)
    total_duration = sum(s.metadata.get("duration_sec", 0) for s in sources)
    headings = sum(s.text.count("\n#") for s in sources if s.source_type in {"text", "md"})
    return {
        "total_chars": total_chars,
        "total_duration_sec": total_duration,
        "heading_estimate": max(1, headings),
        "source_count": len(sources),
    }
