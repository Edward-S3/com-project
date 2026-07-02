"""python-pptx によるレンダリング・自由生成モード。"""

from __future__ import annotations

import ast
import json
import logging
import os
import re
import subprocess
import sys
import tempfile
import traceback
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from config.design_system import (
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    DesignContext,
    make_design_bridge,
    palette_from_name,
)
from core.orchestrator import Orchestrator
from templates.layouts import add_background_frame, apply_layout, create_presentation

logger = logging.getLogger(__name__)

# ASTベースで拒否するトップレベルモジュール / 組み込み呼び出し
FORBIDDEN_IMPORT_ROOTS = frozenset(
    {
        "os",
        "sys",
        "subprocess",
        "socket",
        "requests",
        "urllib",
        "http",
        "shutil",
        "pathlib",
        "importlib",
        "ctypes",
        "builtins",
        "runpy",
        "pickle",
        "sqlite3",
    }
)
FORBIDDEN_CALL_NAMES = frozenset({"eval", "exec", "open", "__import__", "compile", "globals", "locals"})

MAX_CODEGEN_ATTEMPTS = 2


def _build_exec_namespace(slide, data: dict, design: DesignContext) -> Dict[str, Any]:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

    bridge = make_design_bridge(design)
    return {
        "slide": slide,
        "data": data,
        "design": bridge,
        "Inches": Inches,
        "Pt": Pt,
        "RGBColor": RGBColor,
        "MSO_SHAPE": MSO_SHAPE,
        "MSO_ANCHOR": MSO_ANCHOR,
        "MSO_AUTO_SIZE": MSO_AUTO_SIZE,
        "PP_ALIGN": PP_ALIGN,
        "SLIDE_WIDTH": bridge.SLIDE_WIDTH,
        "SLIDE_HEIGHT": bridge.SLIDE_HEIGHT,
    }


def _sandbox_wrapper_script(code: str, slide_data: dict, design: DesignContext, out_pptx: str) -> str:
    data_json = json.dumps(slide_data, ensure_ascii=False)
    return f"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from config.design_system import palette_from_name, make_design_bridge

design = make_design_bridge(palette_from_name({json.dumps(design.palette_name)!r}))
data = json.loads({data_json!r})
SLIDE_WIDTH = design.SLIDE_WIDTH
SLIDE_HEIGHT = design.SLIDE_HEIGHT

{code}

prs = Presentation()
prs.slide_width = Inches({SLIDE_WIDTH_IN})
prs.slide_height = Inches({SLIDE_HEIGHT_IN})
slide = prs.slides.add_slide(prs.slide_layouts[6])
render_slide(slide, data, design)
prs.save({out_pptx!r})
"""


def _is_api_timeout_error(exc: BaseException) -> bool:
    msg = str(exc)
    return any(k in msg for k in ("504", "DEADLINE_EXCEEDED", "タイムアウト", "TimeoutError"))


class SandboxError(Exception):
    """サンドボックス実行失敗。"""

    def __init__(self, stage: str, message: str, code: str) -> None:
        super().__init__(message)
        self.stage = stage
        self.code = code


def _load_timeouts() -> dict:
    path = Path(__file__).resolve().parent.parent / "config" / "timeouts.json"
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def _extract_generated_code(raw: str) -> str:
    """LLM応答から render_slide 定義を抽出。"""
    text = raw.strip()
    fence = re.search(r"```(?:python)?\s*(.*?)```", text, re.S | re.I)
    if fence:
        text = fence.group(1).strip()
    fn = re.search(r"(def\s+render_slide\s*\(.*)", text, re.S)
    if fn:
        text = fn.group(1).strip()
    return text


def _validate_generated_code(code: str) -> Tuple[Optional[str], str]:
    """ASTベースの簡易バリデーション。戻り値: (エラーメッセージ, stage='validation')"""
    try:
        tree = ast.parse(code)
    except SyntaxError as exc:
        return f"構文エラー: {exc}", "validation"

    if not any(isinstance(n, ast.FunctionDef) and n.name == "render_slide" for n in tree.body):
        return "render_slide 関数が定義されていません", "validation"

    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for alias in node.names:
                root = alias.name.split(".")[0]
                if root in FORBIDDEN_IMPORT_ROOTS:
                    return f"禁止された import: {alias.name}", "validation"
        elif isinstance(node, ast.ImportFrom):
            root = (node.module or "").split(".")[0]
            if root in FORBIDDEN_IMPORT_ROOTS:
                return f"禁止された import from: {node.module}", "validation"
        elif isinstance(node, ast.Call) and isinstance(node.func, ast.Name):
            if node.func.id in FORBIDDEN_CALL_NAMES:
                return f"禁止された関数呼び出し: {node.func.id}", "validation"
    return None, "validation"


def _log_sandbox_error(
    log_dir: Path,
    *,
    stage: str,
    exc: BaseException,
    code: str,
    slide_title: str,
    attempt: int,
) -> Path:
    """sandbox_errors_{timestamp}.log に詳細を記録。"""
    log_dir.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = log_dir / f"sandbox_errors_{ts}.log"
    head = "\n".join(code.splitlines()[:20])
    body = (
        f"timestamp={datetime.now().isoformat()}\n"
        f"slide_title={slide_title!r}\n"
        f"attempt={attempt}\n"
        f"failure_stage={stage}\n"
        f"exception_type={type(exc).__name__}\n"
        f"exception_message={exc}\n"
        f"--- traceback ---\n"
        f"{traceback.format_exc()}\n"
        f"--- code (first 20 lines) ---\n"
        f"{head}\n"
        f"--- code (full) ---\n"
        f"{code}\n"
    )
    path.write_text(body, encoding="utf-8")
    logger.error("サンドボックス失敗 stage=%s slide=%r -> %s", stage, slide_title, path)
    return path


def _build_codegen_prompt(slide_data: dict, design: DesignContext, error: Optional[str] = None) -> str:
    err = f"\n前回のエラー: {error}\n" if error else ""
    return f"""python-pptxでスライド1枚を描画する関数を生成してください。

関数シグネチャ(この名前・引数を厳守):
def render_slide(slide, data, design):
    ...

data = {json.dumps(slide_data, ensure_ascii=False)}
【配色ルール — 白背景前提で必ず守ること】
- 背景: design.body_bg() / design.title_bg() のみ（白）
- 見出し・本文: design.title_fg() / design.body_fg() / design.primary
- サブタイトル・補足・キャプション: design.subtitle_fg() / design.text_muted
- 数値強調・リンク風テキスト: design.accent_text()
- design.secondary / design.surface は図形・カードの塗りつぶし専用。文字色に使わない
- 淡いパステル色を文字色に使わない（視認性が低下する）
design.body_bg(), design.body_fg(), design.title_bg(), design.title_fg() は RGBColor を返します。
design.primary, design.accent, design.text_muted, design.subtitle_fg(), design.accent_text() も RGBColor です。
スライド寸法: SLIDE_WIDTH, SLIDE_HEIGHT (Inches) を使用してください。
slide.presentation や slide.slide_width は使わないでください。
fill.fore_color.rgb = design.body_bg() のように RGBColor を直接代入してください。
RGBColor.from_hex / from_string は使わないでください。

許可: pptx, pptx.util, pptx.dml.color, pptx.enum からの import
禁止: os, sys, subprocess, ファイルI/O, ネットワーク, eval/exec/open
禁止: タイトル下アクセントライン、装飾バー
最小余白0.5インチ、視覚要素を1つ以上含める
{err}
def render_slide で始まる Python コードのみ返してください(説明文・マークダウン不可)。
"""


def _run_codegen_subprocess(
    code: str,
    slide_data: dict,
    design: DesignContext,
    *,
    timeout: int = 5,
    log_dir: Path,
    slide_title: str,
    attempt: int,
) -> None:
    """サンドボックスでコードを検証実行。失敗時は SandboxError + ログファイル。"""
    validation_err, stage = _validate_generated_code(code)
    if validation_err:
        exc = SandboxError(stage, validation_err, code)
        _log_sandbox_error(log_dir, stage=stage, exc=exc, code=code, slide_title=slide_title, attempt=attempt)
        raise exc

    with tempfile.TemporaryDirectory() as tmp:
        script = Path(tmp) / "render_slide.py"
        out_pptx = str(Path(tmp) / "out.pptx")
        wrapper = _sandbox_wrapper_script(code, slide_data, design, out_pptx)
        script.write_text(wrapper, encoding="utf-8")
        project_root = str(Path(__file__).resolve().parent.parent)
        env = os.environ.copy()
        env["PYTHONPATH"] = project_root + os.pathsep + env.get("PYTHONPATH", "")
        try:
            result = subprocess.run(
                [sys.executable, str(script)],
                capture_output=True,
                text=True,
                timeout=timeout,
                cwd=project_root,
                env=env,
            )
        except subprocess.TimeoutExpired as exc:
            err = SandboxError("execution", f"サンドボックス実行タイムアウト ({timeout}s)", code)
            _log_sandbox_error(log_dir, stage="execution", exc=err, code=code, slide_title=slide_title, attempt=attempt)
            raise err from exc

        if result.returncode != 0:
            msg = (result.stderr or result.stdout or "不明な実行エラー").strip()
            err = SandboxError("execution", msg, code)
            _log_sandbox_error(log_dir, stage="execution", exc=err, code=code, slide_title=slide_title, attempt=attempt)
            raise err


def _execute_render_code(code: str, slide, data: dict, design: DesignContext) -> None:
    """検証済みコードを本番スライドに適用。"""
    namespace = _build_exec_namespace(slide, data, design)
    exec(compile(code, "<generated_slide>", "exec"), namespace)  # noqa: S102
    render_fn = namespace.get("render_slide")
    if not callable(render_fn):
        raise SandboxError("execution", "render_slide が実行コンテキストに存在しません", code)
    render_fn(slide, data, namespace["design"])


def render_slide_with_freedom(
    orchestrator: Orchestrator,
    prs: Presentation,
    slide_data: dict,
    design: DesignContext,
    *,
    template_only: bool = False,
    use_background_frame: bool = False,
    cancel_event=None,
    log_dir: Path = Path("logs"),
) -> Tuple[str, str]:
    """戻り値: (mode, layout_name) mode=free|fallback|template"""
    if template_only:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if use_background_frame:
            add_background_frame(slide)
        name = apply_layout(slide, slide_data, design, use_background_frame=False)
        return "template", name

    timeouts = _load_timeouts()
    gen_timeout = int(timeouts.get("layout_generation_sec", 30))
    exec_timeout = int(timeouts.get("layout_execution_sec", 5))
    slide_title = slide_data.get("title", "")

    last_error = None
    for attempt in range(MAX_CODEGEN_ATTEMPTS):
        if cancel_event and cancel_event.is_set():
            raise InterruptedError("キャンセル")
        code = ""
        try:
            prompt = _build_codegen_prompt(slide_data, design, last_error)
            resp = orchestrator.run_task(
                "slide_layout_code_generation",
                prompt,
                cancel_event=cancel_event,
                timeout_sec=gen_timeout,
            )
            if not resp:
                break
            code = _extract_generated_code(resp.text)
            _run_codegen_subprocess(
                code,
                slide_data,
                design,
                timeout=exec_timeout,
                log_dir=log_dir,
                slide_title=slide_title,
                attempt=attempt + 1,
            )
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if use_background_frame:
                add_background_frame(slide)
            _execute_render_code(code, slide, slide_data, design)
            hint = slide_data.get("layout_hint", "custom")
            _log_render(log_dir, slide_title, "free", hint)
            return "free", hint
        except SandboxError as exc:
            last_error = f"[{exc.stage}] {exc}"
            logger.warning("自由生成失敗 attempt=%s stage=%s: %s", attempt + 1, exc.stage, exc)
        except Exception as exc:
            last_error = str(exc)
            _log_sandbox_error(
                log_dir,
                stage="orchestration",
                exc=exc,
                code=code,
                slide_title=slide_title,
                attempt=attempt + 1,
            )
            logger.warning("自由生成失敗 attempt=%s: %s", attempt + 1, exc)
            if _is_api_timeout_error(exc):
                logger.warning("APIタイムアウトのため当該スライドはテンプレートへ即フォールバック")
                break

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if use_background_frame:
        add_background_frame(slide)
    name = apply_layout(slide, slide_data, design, use_background_frame=False)
    _log_render(log_dir, slide_title, "fallback", name)
    return "fallback", name


def _log_render(log_dir: Path, title: str, mode: str, layout: str) -> None:
    log_dir.mkdir(parents=True, exist_ok=True)
    line = f"{datetime.now().isoformat()} render title={title!r} mode={mode} layout={layout}\n"
    with open(log_dir / "render.log", "a", encoding="utf-8") as f:
        f.write(line)


def render_presentation(
    slides: List[Dict[str, Any]],
    notes: List[str],
    design: DesignContext,
    orchestrator: Orchestrator,
    *,
    template_only: bool = False,
    use_background_frame: bool = False,
    cancel_event=None,
    output_path: Path,
) -> Presentation:
    prs = create_presentation(design)
    for i, slide_data in enumerate(slides):
        render_slide_with_freedom(
            orchestrator,
            prs,
            slide_data,
            design,
            template_only=template_only,
            use_background_frame=use_background_frame,
            cancel_event=cancel_event,
        )
        if i < len(notes) and notes[i]:
            try:
                slide = prs.slides[i]
                if slide.has_notes_slide:
                    slide.notes_slide.notes_text_frame.text = notes[i]
            except Exception:
                pass
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return prs
