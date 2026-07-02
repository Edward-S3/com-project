"""再利用可能なスライドレイアウト関数群。"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from config.design_system import (
    BACKGROUND_FRAME_PATH,
    BLOCK_GAP_IN,
    DEFAULT_FONT,
    FRAME_MARGIN_IN,
    MIN_MARGIN_IN,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPOGRAPHY,
    DesignContext,
)
from core.frame_image import get_white_background_frame_path


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _send_to_back(slide, shape) -> None:
    """図形をスライド最背面へ移動（nvGrpSpPr / grpSpPr の直後）。"""
    sp_tree = slide.shapes._spTree  # noqa: SLF001
    element = shape._element  # noqa: SLF001
    sp_tree.remove(element)
    # OOXML では spTree の先頭2要素は nvGrpSpPr と grpSpPr が必須。
    # index=0 に挿入すると PowerPoint がファイルを開けなくなる。
    insert_at = min(2, len(sp_tree))
    sp_tree.insert(insert_at, element)


def add_background_frame(
    slide,
    frame_path: Optional[Path] = None,
) -> bool:
    """背景フレーム画像をスライド中央・最背面に配置（四辺 5mm 余白）。"""
    src = frame_path or BACKGROUND_FRAME_PATH
    path = get_white_background_frame_path(src)
    if not path or not path.exists():
        return False

    try:
        from PIL import Image

        with Image.open(path) as img:
            img_w, img_h = img.size
    except Exception:
        img_w, img_h = 1024, 571

    margin = FRAME_MARGIN_IN
    max_w = SLIDE_WIDTH_IN - 2 * margin
    max_h = SLIDE_HEIGHT_IN - 2 * margin
    if max_w <= 0 or max_h <= 0:
        return False

    img_aspect = img_w / img_h
    box_aspect = max_w / max_h
    if img_aspect > box_aspect:
        pic_w = max_w
        pic_h = max_w / img_aspect
    else:
        pic_h = max_h
        pic_w = max_h * img_aspect

    left = (SLIDE_WIDTH_IN - pic_w) / 2
    top = (SLIDE_HEIGHT_IN - pic_h) / 2

    pic = slide.shapes.add_picture(
        str(path),
        Inches(left),
        Inches(top),
        width=Inches(pic_w),
        height=Inches(pic_h),
    )
    _send_to_back(slide, pic)
    return True


def _set_slide_bg(slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_hex)


def _add_textbox(slide, left, top, width, height, text, *, size=15, bold=False, color="000000", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = DEFAULT_FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    p.alignment = align
    return box


def _add_accent_shape(slide, design: DesignContext, left, top, width, height):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(design.surface_fill())
    shape.line.fill.background()
    return shape


def layout_title(slide, data: Dict[str, Any], design: DesignContext) -> None:
    _set_slide_bg(slide, design.title_bg())
    _add_accent_shape(slide, design, MIN_MARGIN_IN, SLIDE_HEIGHT_IN - 1.2, 2.5, 0.12)
    _add_textbox(
        slide,
        MIN_MARGIN_IN,
        2.0,
        SLIDE_WIDTH_IN - MIN_MARGIN_IN * 2,
        2.0,
        data.get("title", ""),
        size=TYPOGRAPHY["slide_title"]["size_pt"],
        bold=True,
        color=design.title_fg(),
    )
    if data.get("subtitle"):
        _add_textbox(
            slide,
            MIN_MARGIN_IN,
            4.2,
            SLIDE_WIDTH_IN - MIN_MARGIN_IN * 2,
            1.5,
            data["subtitle"],
            size=TYPOGRAPHY["body"]["size_pt"],
            color=design.subtitle_fg(),
        )


def layout_two_column(slide, data: Dict[str, Any], design: DesignContext) -> None:
    _set_slide_bg(slide, design.body_bg())
    _add_textbox(
        slide,
        MIN_MARGIN_IN,
        MIN_MARGIN_IN,
        SLIDE_WIDTH_IN - MIN_MARGIN_IN * 2,
        0.8,
        data.get("title", ""),
        size=TYPOGRAPHY["section_heading"]["size_pt"],
        bold=True,
        color=design.body_fg(),
    )
    col_w = (SLIDE_WIDTH_IN - MIN_MARGIN_IN * 2 - BLOCK_GAP_IN) / 2
    y = MIN_MARGIN_IN + 1.0
    cols = data.get("columns") or [data.get("bullets", []), []]
    if len(cols) < 2:
        mid = len(cols[0]) // 2 if cols else 0
        bullets = data.get("bullets", [])
        cols = [bullets[:mid], bullets[mid:]]
    for i, col in enumerate(cols[:2]):
        text = "\n".join(f"• {b}" for b in col) if isinstance(col, list) else str(col)
        _add_textbox(slide, MIN_MARGIN_IN + i * (col_w + BLOCK_GAP_IN), y, col_w, 4.5, text, color=design.body_fg())
    _add_accent_shape(slide, design, SLIDE_WIDTH_IN - 2.0, y, 1.2, 1.2)


def layout_grid(slide, data: Dict[str, Any], design: DesignContext) -> None:
    _set_slide_bg(slide, design.body_bg())
    _add_textbox(
        slide,
        MIN_MARGIN_IN,
        MIN_MARGIN_IN,
        10,
        0.8,
        data.get("title", ""),
        size=TYPOGRAPHY["section_heading"]["size_pt"],
        bold=True,
        color=design.body_fg(),
    )
    bullets = data.get("bullets", [])[:6]
    positions = [(0, 0), (1, 0), (0, 1), (1, 1), (0, 2), (1, 2)]
    w, h = 5.5, 1.5
    for i, b in enumerate(bullets):
        cx, cy = positions[i]
        x = MIN_MARGIN_IN + cx * (w + BLOCK_GAP_IN)
        y = 1.5 + cy * (h + 0.2)
        _add_accent_shape(slide, design, x, y, w, h)
        _add_textbox(slide, x + 0.2, y + 0.2, w - 0.4, h - 0.4, b, size=14, color=design.body_fg())


def layout_stats(slide, data: Dict[str, Any], design: DesignContext) -> None:
    _set_slide_bg(slide, design.body_bg())
    _add_textbox(
        slide,
        MIN_MARGIN_IN,
        MIN_MARGIN_IN,
        10,
        0.8,
        data.get("title", ""),
        size=TYPOGRAPHY["section_heading"]["size_pt"],
        bold=True,
        color=design.body_fg(),
    )
    stats: List[Dict] = data.get("stats") or []
    if not stats and data.get("bullets"):
        stats = [{"label": b, "value": ""} for b in data["bullets"][:3]]
    for i, st in enumerate(stats[:3]):
        x = MIN_MARGIN_IN + i * 4.0
        _add_accent_shape(slide, design, x, 2.0, 3.5, 2.5)
        _add_textbox(slide, x + 0.2, 2.2, 3.1, 1.0, str(st.get("value", "")), size=36, bold=True, color=design.accent_text(), align=PP_ALIGN.CENTER)
        _add_textbox(slide, x + 0.2, 3.5, 3.1, 0.8, str(st.get("label", "")), size=14, color=design.body_fg(), align=PP_ALIGN.CENTER)


def layout_timeline(slide, data: Dict[str, Any], design: DesignContext) -> None:
    _set_slide_bg(slide, design.body_bg())
    _add_textbox(slide, MIN_MARGIN_IN, MIN_MARGIN_IN, 10, 0.8, data.get("title", ""), size=22, bold=True, color=design.body_fg())
    bullets = data.get("bullets", [])[:5]
    for i, b in enumerate(bullets):
        y = 1.5 + i * 1.0
        _add_accent_shape(slide, design, MIN_MARGIN_IN, y + 0.15, 0.3, 0.3)
        _add_textbox(slide, MIN_MARGIN_IN + 0.6, y, 10, 0.8, b, size=15, color=design.body_fg())


def layout_comparison(slide, data: Dict[str, Any], design: DesignContext) -> None:
    layout_two_column(slide, data, design)


def layout_icon_rows(slide, data: Dict[str, Any], design: DesignContext) -> None:
    _set_slide_bg(slide, design.body_bg())
    _add_textbox(slide, MIN_MARGIN_IN, MIN_MARGIN_IN, 10, 0.8, data.get("title", ""), size=22, bold=True, color=design.body_fg())
    for i, b in enumerate(data.get("bullets", [])[:4]):
        y = 1.5 + i * 1.2
        _add_accent_shape(slide, design, MIN_MARGIN_IN, y, 0.5, 0.5)
        _add_textbox(slide, MIN_MARGIN_IN + 0.8, y, 10, 0.8, b, size=15, color=design.body_fg())


def layout_section(slide, data: Dict[str, Any], design: DesignContext) -> None:
    _set_slide_bg(slide, design.title_bg())
    _add_textbox(
        slide,
        MIN_MARGIN_IN,
        2.8,
        SLIDE_WIDTH_IN - MIN_MARGIN_IN * 2,
        1.5,
        data.get("title", ""),
        size=40,
        bold=True,
        color=design.title_fg(),
        align=PP_ALIGN.CENTER,
    )
    _add_accent_shape(slide, design, 5.5, 4.5, 2.3, 0.2)


def layout_bullets(slide, data: Dict[str, Any], design: DesignContext) -> None:
    _set_slide_bg(slide, design.body_bg())
    _add_textbox(slide, MIN_MARGIN_IN, MIN_MARGIN_IN, 10, 0.8, data.get("title", ""), size=22, bold=True, color=design.body_fg())
    text = "\n".join(f"• {b}" for b in data.get("bullets", []))
    _add_textbox(slide, MIN_MARGIN_IN, 1.5, SLIDE_WIDTH_IN - MIN_MARGIN_IN * 2, 5.0, text, size=15, color=design.body_fg())
    _add_accent_shape(slide, design, SLIDE_WIDTH_IN - 2.2, 1.5, 1.5, 4.5)


LAYOUT_MAP = {
    "title": layout_title,
    "two_column": layout_two_column,
    "grid": layout_grid,
    "stats": layout_stats,
    "timeline": layout_timeline,
    "comparison": layout_comparison,
    "icon_rows": layout_icon_rows,
    "section": layout_section,
    "icon": layout_icon_rows,
}


def pick_layout(data: Dict[str, Any]) -> str:
    if data.get("is_title_slide") or data.get("is_closing_slide"):
        return "title"
    hint = data.get("layout_hint", "two_column")
    return hint if hint in LAYOUT_MAP else "two_column"


def apply_layout(
    slide,
    data: Dict[str, Any],
    design: DesignContext,
    *,
    use_background_frame: bool = False,
) -> str:
    if use_background_frame:
        add_background_frame(slide)
    name = pick_layout(data)
    fn = LAYOUT_MAP.get(name, layout_two_column)
    fn(slide, data, design)
    return name


def create_presentation(design: DesignContext) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    return prs
