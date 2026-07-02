"""PPTX から埋め込み画像を抽出"""
from __future__ import annotations

import io
import logging
from dataclasses import dataclass

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from config import DEFAULT_DPI

logger = logging.getLogger(__name__)


@dataclass
class SlideImage:
    slide_index: int
    image: Image.Image
    image_dpi: int


@dataclass
class LoadedPresentation:
    slide_width_emu: int
    slide_height_emu: int
    slides: list[SlideImage]


def _shape_area(shape) -> int:
    try:
        return int(shape.width) * int(shape.height)
    except Exception:
        return 0


def _is_picture_candidate(shape) -> bool:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except Exception:
        pass
    try:
        if not getattr(shape, "has_text_frame", True):
            name = (getattr(shape, "name", "") or "").lower()
            if "picture" in name:
                return True
    except Exception:
        pass
    return False


def _image_from_shape(shape) -> Image.Image | None:
    try:
        blob = shape.image.blob
        return Image.open(io.BytesIO(blob)).convert("RGB")
    except Exception as exc:
        logger.warning("画像 blob の読み込みに失敗: %s", exc)
        return None


def _extract_dpi(image: Image.Image) -> int:
    dpi_info = image.info.get("dpi")
    if isinstance(dpi_info, tuple) and dpi_info[0]:
        return max(int(dpi_info[0]), 1)
    if isinstance(dpi_info, (int, float)) and dpi_info:
        return max(int(dpi_info), 1)
    return DEFAULT_DPI


def _pick_main_picture(slide) -> Image.Image | None:
    candidates: list[tuple[int, Image.Image]] = []
    for shape in slide.shapes:
        if not _is_picture_candidate(shape):
            continue
        img = _image_from_shape(shape)
        if img is None:
            continue
        candidates.append((_shape_area(shape), img))
    if not candidates:
        return None
    candidates.sort(key=lambda item: item[0], reverse=True)
    return candidates[0][1]


def load_slide_images(pptx_path: str) -> LoadedPresentation:
    """入力 PPTX から各スライドのメイン画像を抽出する。"""
    prs = Presentation(pptx_path)
    slide_width_emu = int(prs.slide_width)
    slide_height_emu = int(prs.slide_height)

    slides: list[SlideImage] = []
    for idx, slide in enumerate(prs.slides):
        image = _pick_main_picture(slide)
        if image is None:
            logger.warning("スライド %d: 埋め込み画像が見つかりません（スキップ）", idx)
            continue
        dpi = _extract_dpi(image)
        slides.append(SlideImage(slide_index=idx, image=image, image_dpi=dpi))
        logger.info(
            "スライド %d: 画像抽出完了 (%dx%d, dpi=%d)",
            idx,
            image.width,
            image.height,
            dpi,
        )

    if not slides:
        raise ValueError("入力 PPTX から画像を抽出できませんでした。")

    return LoadedPresentation(
        slide_width_emu=slide_width_emu,
        slide_height_emu=slide_height_emu,
        slides=slides,
    )
