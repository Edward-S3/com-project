"""中間 JSON から編集可能 PPTX を再構築

品質前提: このツールは OCR と画像解析に基づく近似復元であり、
元スライドのフォント・配色・レイアウトを完全に再現するものではありません。
出力 PPTX は「テキスト編集の起点」として使用することを想定しています。
"""
from __future__ import annotations

import logging
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu, Pt

from bbox_utils import sanitize_bbox, to_crop_box
from config import DEFAULT_FONT, DEFAULT_SLIDE_HEIGHT_EMU, DEFAULT_SLIDE_WIDTH_EMU

logger = logging.getLogger(__name__)


def px_to_emu(px: float, image_dpi: int) -> int:
    """ピクセル座標を EMU に変換する。EMU = px × 9525 ÷ (image_dpi ÷ 96)"""
    dpi = max(image_dpi, 1)
    return int(px * 9525 / (dpi / 96))


def _create_masked_background(image: Image.Image, bboxes: list[dict]) -> Image.Image:
    bg = image.copy()
    draw = ImageDraw.Draw(bg)
    for bbox in bboxes:
        box = to_crop_box(bbox, image.width, image.height, pad=2)
        if box is None:
            continue
        x1, y1, x2, y2 = box
        draw.rectangle([x1, y1, x2, y2], fill=(255, 255, 255))
    return bg


def _add_full_slide_image(slide, image: Image.Image, slide_width_emu: int, slide_height_emu: int) -> None:
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        path = tmp.name
        image.save(path, format="PNG")
    try:
        slide.shapes.add_picture(path, 0, 0, width=Emu(slide_width_emu), height=Emu(slide_height_emu))
    finally:
        Path(path).unlink(missing_ok=True)


def _add_background(slide, image: Image.Image, data: dict, slide_width_emu: int, slide_height_emu: int) -> None:
    mask_boxes = []
    for block in data.get("text_blocks") or []:
        mask_boxes.append(block["bbox"])
    for block in data.get("image_blocks") or []:
        mask_boxes.append(block["bbox"])
    bg_image = _create_masked_background(image, mask_boxes) if mask_boxes else image
    _add_full_slide_image(slide, bg_image, slide_width_emu, slide_height_emu)


def _add_text_blocks(slide, data: dict) -> None:
    dpi = int(data.get("image_dpi") or 96)
    for block in data.get("text_blocks") or []:
        bbox = block["bbox"]
        left = px_to_emu(bbox["x"], dpi)
        top = px_to_emu(bbox["y"], dpi)
        width = px_to_emu(bbox["width"], dpi)
        height = px_to_emu(bbox["height"], dpi)
        shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = block["content"]
        font = p.font
        font.name = DEFAULT_FONT
        role = block.get("role", "other")
        size_pt = float(block.get("font_size_pt") or 14)
        if role == "title":
            font.bold = True
            size_pt = max(size_pt, 20.0)
        font.size = Pt(size_pt)


def _add_image_blocks(slide, data: dict) -> None:
    dpi = int(data.get("image_dpi") or 96)
    for block in data.get("image_blocks") or []:
        crop_path = block.get("cropped_path")
        if not crop_path or not Path(crop_path).is_file():
            logger.warning("画像ブロック %s: cropped_path が見つかりません", block.get("id"))
            continue
        bbox = block["bbox"]
        left = px_to_emu(bbox["x"], dpi)
        top = px_to_emu(bbox["y"], dpi)
        width = px_to_emu(bbox["width"], dpi)
        height = px_to_emu(bbox["height"], dpi)
        slide.shapes.add_picture(crop_path, Emu(left), Emu(top), Emu(width), Emu(height))


def _blank_layout(prs: Presentation):
    if prs.slide_layouts:
        return prs.slide_layouts[6]
    return prs.slide_layouts[0]


def write_presentation(
    slide_structures: list[dict],
    source_images: dict[int, Image.Image],
    output_path: str | Path,
    slide_width_emu: int | None = None,
    slide_height_emu: int | None = None,
) -> None:
    """中間 JSON リストから編集可能 PPTX を生成する。"""
    prs = Presentation()
    width = slide_width_emu or DEFAULT_SLIDE_WIDTH_EMU
    height = slide_height_emu or DEFAULT_SLIDE_HEIGHT_EMU
    prs.slide_width = width
    prs.slide_height = height

    layout = _blank_layout(prs)
    for data in slide_structures:
        slide_index = int(data["slide_index"])
        slide = prs.slides.add_slide(layout)
        source = source_images.get(slide_index)
        if source is None:
            logger.warning("スライド %d: ソース画像がありません", slide_index)
            continue

        if data.get("use_full_image"):
            logger.info("スライド %d: フォールバック（元画像をそのまま配置）", slide_index)
            _add_full_slide_image(slide, source, width, height)
            continue

        _add_background(slide, source, data, width, height)
        _add_image_blocks(slide, data)
        _add_text_blocks(slide, data)
        logger.info(
            "スライド %d: PPTX 配置完了 (text=%d, images=%d)",
            slide_index,
            len(data.get("text_blocks") or []),
            len(data.get("image_blocks") or []),
        )

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    logger.info("PPTX 出力完了: %s", out)
